"""
Gerador de slides — Intensivao Quant AI | Aulas 2 a 10
Execucao: python gerar_slides_aulas02_10.py
Saida: um .pptx por aula nas respectivas pastas
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0D, 0x1B, 0x3E)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
GRAY  = RGBColor(0x77, 0x88, 0x9A)
GREEN = RGBColor(0x1E, 0x8B, 0x4C)
BLUE  = RGBColor(0x1A, 0x6E, 0xAE)
RED   = RGBColor(0xC0, 0x39, 0x2B)
PURP  = RGBColor(0x6C, 0x35, 0x83)
ORAN  = RGBColor(0xBF, 0x6A, 0x02)

W, H = 13.33, 7.5
BASE = os.path.dirname(os.path.abspath(__file__))


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp

def txt(slide, l, t, w, h, content, size=13, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def buls(slide, l, t, w, h, items, size=13, color=WHITE, spacing=7):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        r = p.add_run()
        r.text = f"  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def ref(slide, citation):
    txt(slide, 0.3, 7.05, W - 0.5, 0.38,
        f"Ref.: {citation}", size=9, color=GRAY, italic=True)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, 1.15, fill=NAVY)
    rect(slide, 0, 0, 0.22, 1.15, fill=GOLD)
    txt(slide, 0.42, 0.1, W - 0.55, 0.58, title,
        size=22, bold=True, color=WHITE)
    if subtitle:
        txt(slide, 0.42, 0.7, W - 0.55, 0.38, subtitle,
            size=12, color=LIGHT, italic=True)

def capa(prs, num, titulo, subtitulo, temas):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=NAVY)
    rect(s, 0, 0, 0.45, H, fill=GOLD)
    rect(s, 0.45, 3.5, W - 0.45, 0.04, fill=GOLD)
    txt(s, 0.85, 0.3, 4.5, 1.1, f"Aula {num}", size=50, bold=True, color=GOLD)
    txt(s, 0.85, 1.5, 11.0, 0.65, titulo, size=26, bold=True, color=WHITE)
    txt(s, 0.85, 2.2, 11.0, 0.42, subtitulo, size=14, color=LIGHT, italic=True)
    txt(s, 0.85, 2.72, 11.0, 0.32,
        "Intensivao Quant AI | ImpactUFSCar | 2025", size=11, color=GRAY)
    txt(s, 0.85, 3.7, 11.0, 0.35, "Nesta aula:", size=13, bold=True, color=GOLD)
    buls(s, 0.85, 4.1, 11.0, 3.0, temas, size=13, color=LIGHT, spacing=6)

def agenda(prs, teoria, codigo):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Estrutura da Aula", "~30 min de teoria + ~90 min de live coding")
    rect(s, 0.4, 1.25, 6.1, 5.55, fill=LIGHT)
    rect(s, 0.4, 1.25, 6.1, 0.42, fill=NAVY)
    txt(s, 0.6, 1.3, 5.8, 0.36, "Fundamentos (teoria)", size=13, bold=True, color=WHITE)
    buls(s, 0.6, 1.78, 5.8, 4.5, teoria, size=13, color=NAVY, spacing=10)
    rect(s, 6.8, 1.25, 6.1, 5.55, fill=LIGHT)
    rect(s, 6.8, 1.25, 6.1, 0.42, fill=GOLD)
    txt(s, 7.0, 1.3, 5.8, 0.36, "Live Coding", size=13, bold=True, color=NAVY)
    buls(s, 7.0, 1.78, 5.8, 4.5, codigo, size=13, color=NAVY, spacing=10)

def construir(prs, funcoes, entradas, saidas):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "O que Vamos Construir Hoje",
           "Criterio de conclusao: notebook roda do inicio ao fim sem erros e gera os parquets abaixo")
    rect(s, 0.4, 1.25, 8.1, 5.55, fill=LIGHT)
    rect(s, 0.4, 1.25, 8.1, 0.42, fill=NAVY)
    txt(s, 0.6, 1.3, 7.8, 0.36, "Funcoes e logica que vamos implementar:", size=13, bold=True, color=WHITE)
    buls(s, 0.6, 1.78, 7.8, 4.7, funcoes, size=13, color=NAVY, spacing=9)
    rect(s, 8.75, 1.25, 4.2, 2.6, fill=RGBColor(0xEA, 0xF4, 0xFE))
    txt(s, 8.95, 1.3, 3.9, 0.36, "Entradas:", size=12, bold=True, color=NAVY)
    buls(s, 8.95, 1.72, 3.9, 1.95, entradas, size=11, color=NAVY, spacing=6)
    rect(s, 8.75, 4.1, 4.2, 2.7, fill=RGBColor(0xEA, 0xF8, 0xEE))
    txt(s, 8.95, 4.15, 3.9, 0.36, "Saidas (parquets):", size=12, bold=True, color=NAVY)
    buls(s, 8.95, 4.57, 3.9, 2.1, saidas, size=11, color=NAVY, spacing=6)

def card2(slide, l, t, w, h, title, title_bg, body, size=12):
    rect(slide, l, t, w, h, fill=LIGHT)
    rect(slide, l, t, w, 0.42, fill=title_bg)
    txt(slide, l + 0.12, t + 0.07, w - 0.2, 0.34, title, size=12, bold=True, color=WHITE)
    buls(slide, l + 0.12, t + 0.52, w - 0.2, h - 0.62, body, size=size, color=NAVY, spacing=5)


# ══════════════════════════════════════════════════════════
# AULA 02 — DADOS
# ══════════════════════════════════════════════════════════
def gerar_aula02():
    prs = new_prs()
    capa(prs, "02", "Dados: Coleta, Limpeza e Pipeline",
         "A fundacao de dados que todas as aulas seguintes vao consumir",
         ["O IBOVESPA: composicao, criterios de inclusao e ~77 acoes",
          "yfinance: como funciona a coleta automatica de precos historicos",
          "Qualidade de dados: missing values, splits, dividendos e outliers",
          "Parquet: por que e melhor que CSV para series temporais",
          "Pipeline reprodutivel: coletar uma vez, usar em todas as aulas"])

    agenda(prs,
        ["O que e o IBOVESPA e como e composto",
         "Como o yfinance acessa dados do Yahoo Finance",
         "Tipos de problemas em dados financeiros brutos",
         "Por que parquet e o formato certo aqui"],
        ["Download de precos de ~77 acoes (.SA)",
         "Calcular retornos diarios (pct_change)",
         "Agregar para retornos mensais (resample ME)",
         "Tratar NaN, outliers e ajustes corporativos",
         "Salvar 3 parquets para as proximas aulas"])

    # Slide 3: O IBOVESPA
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "O IBOVESPA — Nosso Universo de Ativos",
           "Principal indice de acoes do Brasil: criterios de inclusao e implicacoes para o backtest")
    txt(s, 0.5, 1.28, 12.3, 0.48,
        "O IBOVESPA e o indice de referencia do mercado acionario brasileiro, composto pelas acoes "
        "mais liquidas da B3. Sua composicao e revista a cada 4 meses com base em liquidez e presenca "
        "em pregao. Usamos a composicao atual como proxy do universo historico.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.9, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 2.05, 3.98, 4.72, "Criterios de Inclusao", BLUE,
          ["Negociado em >= 95% dos pregoes do periodo de apuracao",
           "Participacao minima em volume financeiro (threshold definido pela B3)",
           "Preco > R$ 1,00 (nao pode ser penny stock)",
           "Nao pode estar em processo de falencia ou recuperacao judicial",
           "Resultado: ~70-85 acoes dependendo do periodo"])
    card2(s, 4.67, 2.05, 3.98, 4.72, "Nossa Abordagem", NAVY,
          ["Lista de ~77 tickers com sufixo .SA (ex: PETR4.SA, VALE3.SA)",
           "Survivorship bias: acoes que saem do indice sao ignoradas",
           "Mitigacao pratica: periodo fixo com tickers historicamente relevantes",
           "Alternativa profissional: usar base point-in-time (Compustat, Refinitiv)",
           "Para o contexto do Desafio: abordagem simplificada e aceitavel"])
    card2(s, 8.95, 2.05, 3.98, 4.72, "Download via yfinance", GREEN,
          ["yfinance.download(tickers, start, end, auto_adjust=True)",
           "Retorna Adj Close (ajustado por splits e dividendos)",
           "Colunas: ticker; Indice: datetime",
           "Periodo recomendado: 2012-01-01 ate hoje",
           "Cuidado: rate limiting — adicionar delays entre downloads"])
    ref(s, "B3 (2024). Metodologia do Indice BOVESPA. b3.com.br")

    # Slide 4: Qualidade de dados
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Qualidade de Dados Financeiros",
           "Cada problema nao tratado e uma fonte de erro que pode invalidar o backtest")
    txt(s, 0.5, 1.28, 12.3, 0.45,
        "Dados financeiros brutos raramente estao prontos para uso direto. "
        "Garbage in, garbage out: um backtest construido sobre dados ruins produz "
        "resultados que nao se replicam fora da amostra.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.87, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 2.02, 3.98, 4.75, "Missing Values", RED,
          ["Causa: ativo suspenso, feriado local, falha de API",
           "Tratamento: .ffill() para falhas curtas (<=5 dias)",
           "Remover acao se ausencia > 20% do periodo total",
           "Nunca interpolar retornos: cria autocorrelacao artificial",
           "Check: df.isnull().sum() / len(df) por coluna"])
    card2(s, 4.67, 2.02, 3.98, 4.75, "Ajustes Corporativos", BLUE,
          ["Split: 1 acao vira N — preco cai mas retorno nao muda",
           "Dividendo: reduz preco na data ex (nao e perda real)",
           "Solucao: usar auto_adjust=True no yfinance",
           "Adj Close incorpora todos os ajustes historicos",
           "Validar: retorno acumulado deve ser consistente"])
    card2(s, 8.95, 2.02, 3.98, 4.75, "Outliers e Erros", ORAN,
          ["Retornos de +200% ou -80% em um mes: quase sempre erro",
           "Deteccao: filtrar |retorno_mensal| > 50%",
           "Validar manualmente cada outlier identificado",
           "Remover o dado, nao a acao inteira (se for pontual)",
           "Registrar todas as limpezas feitas (log de decisoes)"])
    ref(s, "Lopez de Prado, M. (2018). Advances in Financial Machine Learning, Cap. 2. Wiley")

    # Slide 5: Pipeline
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Pipeline de Dados — Arquitetura do Projeto",
           "Rodar uma vez, reutilizar em todas as aulas — reproducibilidade e velocidade")
    steps = [
        ("1. Download", NAVY, "yfinance.download()\n~77 tickers x ~3.000 dias"),
        ("2. Limpeza",  BLUE, "NaN, outliers\nAdj Close verificado"),
        ("3. Retornos", GREEN,"ret_diarios = pct_change()\nret_mensais = resample('ME').last()"),
        ("4. Parquet",  GOLD, "to_parquet(snappy)\nLeitura em <1s nas proximas aulas"),
    ]
    for i, (title, color, body) in enumerate(steps):
        x = 0.4 + i * 3.22
        rect(s, x, 1.9, 3.0, 2.6, fill=color)
        txt(s, x + 0.15, 1.97, 2.7, 0.42, title, size=14, bold=True, color=WHITE)
        txt(s, x + 0.15, 2.45, 2.7, 1.9, body, size=12, color=WHITE)
        if i < 3:
            txt(s, x + 3.05, 2.95, 0.3, 0.5, ">", size=24, bold=True, color=NAVY)
    rect(s, 0.4, 4.65, 12.5, 0.03, fill=GOLD)
    txt(s, 0.5, 4.78, 12.0, 0.35, "Por que Parquet e nao CSV?", size=14, bold=True, color=NAVY)
    buls(s, 0.5, 5.18, 5.9, 1.8,
         ["Leitura 10-50x mais rapida para DataFrames grandes",
          "Preserva tipos (datetime64, float64) automaticamente",
          "Compressao snappy: ~5x menor que CSV sem perda"],
         size=13, color=NAVY, spacing=7)
    buls(s, 6.6, 5.18, 6.3, 1.8,
         ["Suporta leitura parcial (colunas) sem carregar tudo",
          "Padrao do ecossistema (Spark, Polars, DuckDB, Arrow)",
          "API: df.to_parquet('path.parquet') / pd.read_parquet('path.parquet')"],
         size=13, color=NAVY, spacing=7)

    construir(prs,
        ["baixar_precos(tickers, start, end) — download com retry e tratamento de erros",
         "limpar_precos(df_raw) — remover NaN excessivos e outliers extremos",
         "calcular_retornos_diarios(precos) — pct_change() com validacao",
         "calcular_retornos_mensais(ret_diarios) — resample('ME').apply(composto)",
         "pipeline completo: executar celulas em ordem, verificar shape e dtypes ao final"],
        ["(nenhum — esta aula gera os dados base)"],
        ["dados/precos_ibov.parquet",
         "dados/retornos_diarios_limpo.parquet",
         "dados/retornos_mensais_limpo.parquet"])

    out = os.path.join(BASE, "aula-02-dados", "slides-aula-02-dados.pptx")
    prs.save(out)
    print(f"  Aula 02: {len(prs.slides)} slides -> aula-02-dados/")


# ══════════════════════════════════════════════════════════
# AULA 03 — EDA
# ══════════════════════════════════════════════════════════
def gerar_aula03():
    prs = new_prs()
    capa(prs, "03", "EDA: Analise Exploratoria dos Retornos",
         "Entender os dados antes de qualquer estrategia — nunca pule esta etapa",
         ["Por que EDA e obrigatoria antes de qualquer backtest",
          "Distribuicao de retornos: fat tails e o problema com a gaussiana",
          "Correlacoes entre acoes: o que realmente diversifica",
          "Estacionariedade: por que usamos retornos e nao precos",
          "Visualizacoes que revelam padroes (e armadilhas) nos dados"])

    agenda(prs,
        ["Fat tails: por que a gaussiana falha em financas",
         "Correlacoes dinamicas e seu significado",
         "Estacionariedade e o teste ADF",
         "O que o EDA revela sobre o IBOVESPA"],
        ["Distribuicao dos retornos (histograma + QQ-plot)",
         "Heatmap de correlacoes entre acoes",
         "Retorno medio vs volatilidade (scatter)",
         "Evolucao temporal da correlacao media",
         "Estatisticas descritivas completas"])

    # Slide 3: Distribuicao de retornos
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Distribuicao de Retornos — Fat Tails",
           "Por que a hipotese de normalidade falha em dados financeiros reais")
    txt(s, 0.5, 1.28, 12.3, 0.48,
        "O modelo de Black-Scholes e o CAPM assumem retornos normalmente distribuidos. "
        "Na realidade, retornos financeiros apresentam caudas mais pesadas que a gaussiana (fat tails), "
        "assimetria negativa e curtose elevada — o que subestima o risco de eventos extremos.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.9, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 2.05, 3.98, 4.72, "O que sao Fat Tails?", BLUE,
          ["Distribuicao normal: P(|r| > 3σ) ≈ 0,3% dos dias",
           "Mercado real: crashes de 5-10σ ocorrem muito mais",
           "Exemplo: Black Monday 1987 = retorno de -22σ pela gaussiana",
           "Curtose: mede o 'peso' das caudas (normal = 3; mercado > 5)",
           "Skewness negativa: quedas sao maiores que altas (assimetria)"])
    card2(s, 4.67, 2.05, 3.98, 4.72, "Implicacoes Praticas", RED,
          ["VaR gaussiano SUBESTIMA risco de cauda (Basel III reconhece isso)",
           "Sharpe ratio assume normalidade — usar Sortino como complemento",
           "Max drawdown captura o que o desvio-padrao nao captura",
           "Testes estatisticos (t-test) perdem poder com fat tails",
           "Usar bootstrap para intervalos de confianca mais robustos"])
    card2(s, 8.95, 2.05, 3.98, 4.72, "Como Medir no Python", GREEN,
          ["scipy.stats.kurtosis(retornos, fisher=True)  # excesso sobre 3",
           "scipy.stats.skew(retornos)  # negativo = cauda esquerda pesada",
           "scipy.stats.jarque_bera(retornos)  # teste de normalidade",
           "scipy.stats.probplot(retornos)  # QQ-plot visual",
           "retornos.describe(percentiles=[.01,.05,.95,.99])"])
    ref(s, "Mandelbrot, B. (1963). The Variation of Certain Speculative Prices. J. of Business, 36(4), 394–419; "
           "Taleb, N.N. (2007). The Black Swan. Random House")

    # Slide 4: Correlacoes e estacionariedade
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Correlacoes e Estacionariedade",
           "O que realmente diversifica e por que retornos (nao precos) sao o dado certo")
    rect(s, 0.5, 1.25, 12.3, 0.03, fill=GOLD)
    rect(s, 0.4, 1.35, 6.1, 5.45, fill=LIGHT)
    txt(s, 0.6, 1.42, 5.8, 0.4, "Correlacoes entre Acoes", size=14, bold=True, color=NAVY)
    buls(s, 0.6, 1.9, 5.8, 4.2,
         ["Correlacao media do IBOVESPA: ~0,45 em periodos normais",
          "Em crises: correlacoes sobem para ~0,80 (tudo cai junto)",
          "Implicacao: diversificacao falha exatamente quando mais e necessaria",
          "Clusters setoriais: bancos correlacionados entre si, mineracao idem",
          "Heatmap de correlacao: ferramenta visual essencial pre-backtest",
          "Correlacoes nao sao estaveis no tempo — monitorar janelas moveis"],
         size=12, color=NAVY, spacing=7)
    rect(s, 6.8, 1.35, 6.1, 5.45, fill=LIGHT)
    txt(s, 7.0, 1.42, 5.8, 0.4, "Estacionariedade — Por que Retornos?", size=14, bold=True, color=NAVY)
    buls(s, 7.0, 1.9, 5.8, 4.2,
         ["Preco e nao-estacionario: media e variancia mudam no tempo",
          "Preco: tendencia de alta estrutural -> regressao espuria",
          "Retorno = ln(P_t/P_{t-1}): aproximadamente estacionario",
          "Teste ADF (Augmented Dickey-Fuller): p<0.05 rejeita raiz unitaria",
          "Retornos passam no ADF; precos falham consistentemente",
          "Regra: NUNCA modele precos diretamente — sempre retornos"],
         size=12, color=NAVY, spacing=7)
    ref(s, "Dickey, D.A. & Fuller, W.A. (1979). Distribution of Estimators for AR Time Series with a Unit Root. "
           "JASA, 74(366), 427–431")

    construir(prs,
        ["distribuicao_retornos(ret) — histograma + curva normal + estatisticas (kurt, skew)",
         "qq_plot(ret) — quantil-quantil para visualizar desvio da normalidade",
         "heatmap_correlacao(ret_mensais) — correlacoes medias do periodo todo",
         "scatter_risco_retorno(ret_mensais) — retorno medio x volatilidade por acao",
         "correlacao_movel(ret_mensais, janela=12) — evolucao temporal da correlacao media"],
        ["dados/retornos_diarios_limpo.parquet",
         "dados/retornos_mensais_limpo.parquet"],
        ["(EDA nao salva parquets — saida e visualizacoes e insights)"])

    out = os.path.join(BASE, "aula-03-eda", "slides-aula-03-eda.pptx")
    prs.save(out)
    print(f"  Aula 03: {len(prs.slides)} slides -> aula-03-eda/")


# ══════════════════════════════════════════════════════════
# AULA 04 — SINAL v1
# ══════════════════════════════════════════════════════════
def gerar_aula04():
    prs = new_prs()
    capa(prs, "04", "Sinal v1: Momentum Cross-Sectional",
         "Transformando a hipotese academica em codigo — o fator 12-1 no IBOVESPA",
         ["O que e um fator quantitativo e como o sinal se transforma em posicao",
          "O sinal 12-1: matematica, intuicao e por que o skip de 1 mes importa",
          "Cross-sectional ranking: como ordenar 77 acoes em cada mes",
          "Information Coefficient (IC): como medir se o sinal tem poder preditivo",
          "Diagnostico do sinal antes de qualquer backtest"])

    agenda(prs,
        ["O que e um fator quantitativo (sinal → posicao)",
         "Matematica do sinal 12-1 (com e sem skip)",
         "Cross-sectional vs time-series momentum",
         "IC: a metrica de qualidade do sinal"],
        ["Calcular sinal 12-1 para todas as acoes",
         "Implementar cross-sectional ranking mensal",
         "Computar IC (correlacao sinal x retorno futuro)",
         "Visualizar distribuicao do sinal no tempo",
         "Identificar as acoes mais frequentes no top 10"])

    # Slide 3: O que e um fator
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "O que e um Fator Quantitativo?",
           "Da hipotese economica ao sinal computavel — o pipeline de um fator")
    txt(s, 0.5, 1.28, 12.3, 0.45,
        "Um fator quantitativo e uma caracteristica mensuravel de um ativo que, segundo hipotese "
        "economica fundada em evidencia empirica, prediz seu retorno futuro. "
        "O fator transforma informacao observavel em sinal de compra/venda.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.87, 12.3, 0.03, fill=GOLD)
    steps2 = [
        ("Hipotese\nEconomica", NAVY,
         "Momentum existe pois investidores underreact a noticias (Hong & Stein, 1999). "
         "Acoes com bom retorno recente tem retorno futuro superior em media."),
        ("Sinal\nComputavel", BLUE,
         "ret_mensais.shift(2).rolling(11).sum()\nRetorno acumulado 11 meses\n"
         "com skip de 1 mes. Calculado para cada acao em cada data."),
        ("Ranking\nCross-Sectional", GREEN,
         "Em cada mes, ordenar as ~77 acoes pelo sinal. "
         "Top 10: comprar. Restante: ignorar (long-only). "
         "Sem posicao vendida nesta versao."),
        ("Posicao\nno Portfolio", GOLD,
         "Cada acao do top 10 recebe peso 1/10 (equal-weight nesta versao). "
         "Portfolio rebalanceado mensalmente no fechamento."),
    ]
    for i, (title, color, body) in enumerate(steps2):
        x = 0.4 + i * 3.22
        rect(s, x, 2.02, 3.05, 4.78, fill=LIGHT)
        rect(s, x, 2.02, 3.05, 0.6, fill=color)
        txt(s, x + 0.12, 2.07, 2.82, 0.52, title,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, x + 0.12, 2.7, 2.82, 4.0, body, size=11, color=NAVY)
        if i < 3:
            txt(s, x + 3.1, 4.1, 0.25, 0.5, ">", size=22, bold=True, color=NAVY)
    ref(s, "Hong, H. & Stein, J. (1999). J. of Finance, 54(6), 2143–2184; "
           "Jegadeesh & Titman (1993). J. of Finance, 48(1), 65–91")

    # Slide 4: O sinal 12-1
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "O Sinal 12-1 — Matematica e Intuicao",
           "Por que 12 meses de formacao? Por que pular o ultimo mes?")
    rect(s, 0.5, 1.25, 12.3, 0.03, fill=GOLD)
    rect(s, 0.4, 1.35, 6.1, 5.45, fill=LIGHT)
    txt(s, 0.6, 1.42, 5.8, 0.38, "O Sinal em Python", size=14, bold=True, color=NAVY)
    txt(s, 0.6, 1.88, 5.8, 2.5,
        "# Retorno acumulado de t-12 ate t-2 (skip 1 mes)\nsinal_v1 = ret_mensais.shift(2).rolling(11).sum()\n\n"
        "# Sem o skip (ERRADO — inclui reversao de curto prazo):\nsinal_errado = ret_mensais.shift(1).rolling(12).sum()\n\n"
        "# O .shift(2) garante que na data t usamos\n# informacao disponivel ate t-2 (evita look-ahead)",
        size=11, color=NAVY)
    rect(s, 0.6, 4.5, 5.8, 0.03, fill=GRAY)
    buls(s, 0.6, 4.62, 5.8, 2.1,
         ["shift(2): exclui mes t-1 (skip) E mes t (look-ahead)",
          "rolling(11): janela de 11 meses = sinal 12-1",
          "Resultado: sinal formado entre t-12 e t-2"],
         size=12, color=NAVY, spacing=7)
    rect(s, 6.8, 1.35, 6.1, 5.45, fill=LIGHT)
    txt(s, 7.0, 1.42, 5.8, 0.38, "Por que Skip de 1 Mes?", size=14, bold=True, color=NAVY)
    buls(s, 7.0, 1.88, 5.8, 4.7,
         ["Reversao de curto prazo: o ultimo mes frequentemente reverte (Jegadeesh, 1990)",
          "Causa: market microstructure — bid-ask bounce, liquidez momentanea",
          "Se incluir mes t-1: contamina o sinal com ruido de microestrutura",
          "Evidencia: sinal 12-1 tem IC maior que sinal 12-0 na maioria dos mercados",
          "Padrao da literatura: quase todos os estudos de momentum usam skip-1",
          "No Brasil: reversao de 1 mes e especialmente forte em acoes pequenas"],
         size=12, color=NAVY, spacing=7)
    ref(s, "Jegadeesh, N. (1990). Evidence of Predictable Behavior of Security Returns. "
           "J. of Finance, 45(3), 881–898")

    # Slide 5: IC
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Information Coefficient (IC) — Qualidade do Sinal",
           "Antes de backtester, medir se o sinal tem poder preditivo real")
    txt(s, 0.5, 1.28, 12.3, 0.45,
        "O Information Coefficient e a correlacao de Spearman entre o ranking do sinal em t "
        "e o ranking do retorno realizado em t+1. Mede o poder preditivo do sinal "
        "independentemente da estrategia de portfolio — um bom sinal tem IC positivo e estavel.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.87, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 2.02, 3.98, 4.75, "Calculo do IC", NAVY,
          ["IC_t = spearmanr(sinal_t, retorno_{t+1})[0]",
           "Spearman: correlacao de rankings (robusta a outliers)",
           "IC positivo: sinal prediz retorno na direcao certa",
           "IC medio tipico para momentum: 0,03 a 0,08",
           "IC de 0,05 ja e considerado util na pratica quant"])
    card2(s, 4.67, 2.02, 3.98, 4.75, "Interpretacao", BLUE,
          ["IC > 0: sinal tem valor preditivo (direcao certa)",
           "IC medio > 0,03: sinal pode ser exploravel",
           "IC t-stat = IC_medio / (IC_std / sqrt(N))",
           "t-stat > 2: IC estatisticamente significativo",
           "Hit rate = % de meses em que IC > 0"])
    card2(s, 8.95, 2.02, 3.98, 4.75, "Limitacoes do IC", RED,
          ["IC positivo nao garante alpha apos custos",
           "IC pode variar muito no tempo (regime dependente)",
           "Autocorrelacao do sinal afeta o t-stat naive",
           "Complementar com backtest completo sempre",
           "Decay do IC: medir IC em t+1, t+2, t+3"])
    ref(s, "Grinold, R. & Kahn, R. (1999). Active Portfolio Management, 2a ed. McGraw-Hill — Cap. 6 (IC e IR)")

    construir(prs,
        ["calcular_sinal_v1(ret_mensais) — ret_mensais.shift(2).rolling(11).sum()",
         "ranking_crosssectional(sinal) — .rank(axis=1, pct=True) em cada mes",
         "calcular_ic(sinal, ret_mensais) — spearmanr(sinal_t, ret_{t+1}) para cada t",
         "plot_ic_serie(ic_serie) — IC ao longo do tempo com media e IC_t-stat",
         "top10_frequencia(sinal, n=10) — quais acoes aparecem mais no top 10"],
        ["dados/retornos_mensais_limpo.parquet"],
        ["dados/sinal_v1.parquet"])

    out = os.path.join(BASE, "aula-04-sinal-v1", "slides-aula-04-sinal-v1.pptx")
    prs.save(out)
    print(f"  Aula 04: {len(prs.slides)} slides -> aula-04-sinal-v1/")


# ══════════════════════════════════════════════════════════
# AULA 05 — BACKTEST v1
# ══════════════════════════════════════════════════════════
def gerar_aula05():
    prs = new_prs()
    capa(prs, "05", "Backtest v1: Simulacao Historica",
         "Transformar o sinal em retornos de portfolio e medir a performance",
         ["O que e um backtest e quais hipoteses simplificadoras fazemos",
          "Construcao do portfolio: pesos, retornos e rebalanceamento mensal",
          "Metricas de performance: Sharpe, Sortino, Max Drawdown, Calmar",
          "Benchmark: por que comparar com o IBOVESPA e como fazer certo",
          "Primeiras conclusoes: o sinal gera alpha? A que custo?"])

    agenda(prs,
        ["Hipoteses do backtest: o que assumimos e o que ignoramos",
         "Metricas de risco-retorno: Sharpe, Sortino, MDD, Calmar",
         "Benchmark e alpha: o que significa bater o indice",
         "Limitacoes desta versao (a corrigir na Aula 8)"],
        ["Construir portfolio com top 10 acoes (equal-weight)",
         "Calcular retornos mensais do portfolio",
         "Implementar calcular_metricas() reutilizavel",
         "Comparar com IBOVESPA (benchmark)",
         "Grafico de retorno acumulado e drawdown"])

    # Slide 3: Hipoteses do backtest
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Hipoteses do Backtest — O que Assumimos",
           "Todo backtest e uma simplificacao da realidade — conhecer as hipoteses e fundamental")
    txt(s, 0.5, 1.28, 12.3, 0.45,
        "Um backtest simula como a estrategia teria se comportado no passado. "
        "A distancia entre o backtest e a realidade e determinada pelas hipoteses assumidas. "
        "Hipoteses otimistas geram resultados que nao se replicam fora da amostra.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.87, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 2.02, 3.98, 4.75, "Assumimos (Simplificacoes)", GREEN,
          ["Execucao no fechamento do mes sem slippage",
           "Liquidez ilimitada — qualquer tamanho de posicao",
           "Sem custos de transacao (corrigido na Aula 8)",
           "Dados historicos disponiveis ex-ante (cuidado!)",
           "Portfolio long-only sem alavancagem"])
    card2(s, 4.67, 2.02, 3.98, 4.75, "Ignoramos (Problemas Reais)", RED,
          ["Bid-ask spread: custo real de executar ordens",
           "Market impact: ordens grandes movem o preco",
           "Borrow cost: custo de aluguel em vendas a descoberto",
           "Impostos: IOF, IR sobre ganhos de capital",
           "Restricoes de short: nem toda acao tem aluguel"])
    card2(s, 8.95, 2.02, 3.98, 4.75, "Quando o Backtest e Valido?", BLUE,
          ["Regra geral: resultado bom no backtest simples e necessario, nao suficiente",
           "Validacao fora da amostra (walk-forward) e obrigatoria",
           "Quanto maior o Sharpe simulado, maior o ceticismo necessario",
           "Transparencia total sobre hipoteses ao apresentar resultados",
           "Backtest e ferramenta de falsificacao, nao de confirmacao"])
    ref(s, "Bailey, D.H. & Lopez de Prado, M. (2012). The Sharpe Ratio Efficient Frontier. "
           "JOIS, 5(2), 13–31")

    # Slide 4: Metricas
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Metricas de Performance — O Dicionario do Backtest",
           "Saber o que cada metrica mede e crucial para nao se enganar com bons numeros")
    rect(s, 0.5, 1.25, 12.3, 0.03, fill=GOLD)
    metricas = [
        ("Sharpe Ratio", NAVY,
         "Formula: (E[r] - rf) / σ × √12\nMede: retorno por unidade de risco total\n"
         "Bom: > 0,5 (aceitavel), > 1,0 (bom), > 1,5 (excelente)\n"
         "Limitacao: penaliza retornos positivos acima da media"),
        ("Sortino Ratio", BLUE,
         "Formula: (E[r] - rf) / downside_vol × √12\nMede: retorno por unidade de risco de queda\n"
         "Downside vol: desvio dos retornos negativos apenas\n"
         "Mais relevante para estrategias com retornos assimetricos"),
        ("Max Drawdown", RED,
         "MDD = min( (P_t - max(P_s, s≤t)) / max(P_s, s≤t) )\nMede: pior queda pico-a-vale no periodo\n"
         "Critico para aversao a perda de investidores reais\n"
         "Complementar Sharpe: estrategia pode ter bom Sharpe e MDD devastador"),
        ("Calmar Ratio", GREEN,
         "Formula: CAGR / |Max Drawdown|\nMede: retorno anualizado por unidade de drawdown maximo\n"
         "Bom: > 0,5; estrategias com Calmar > 1,0 sao raras e atrativas\n"
         "Muito valorizado por gestores de risco e alocadores"),
    ]
    for i, (name, color, body) in enumerate(metricas):
        row, col = i // 2, i % 2
        x = 0.4 + col * 6.45
        y = 1.38 + row * 2.82
        rect(s, x, y, 6.2, 2.65, fill=LIGHT)
        rect(s, x, y, 6.2, 0.4, fill=color)
        txt(s, x + 0.15, y + 0.06, 5.9, 0.32, name, size=13, bold=True, color=WHITE)
        txt(s, x + 0.15, y + 0.5, 5.9, 2.0, body, size=11, color=NAVY)
    ref(s, "Sharpe, W.F. (1994). The Sharpe Ratio. J. of Portfolio Management, 21(1), 49–58; "
           "Sortino, F.A. & Price, L.N. (1994). Performance Measurement in a Downside Risk Framework. JPM")

    construir(prs,
        ["construir_portfolio(sinal, ret_mensais, n_top=10) — pesos e retornos mensais",
         "calcular_metricas(retornos, benchmark, nome) — Sharpe, Sortino, MDD, Calmar, alpha, beta",
         "exibir_metricas(*dicts) — tabela comparativa formatada",
         "plot_retorno_acumulado(retornos_dict) — curvas de riqueza comparativas",
         "plot_drawdown(retornos) — evolucao do drawdown ao longo do tempo"],
        ["dados/retornos_mensais_limpo.parquet",
         "dados/sinal_v1.parquet"],
        ["dados/pesos_v1.parquet",
         "dados/retorno_carteira_v1.parquet"])

    out = os.path.join(BASE, "aula-05-backtest-v1", "slides-aula-05-backtest-v1.pptx")
    prs.save(out)
    print(f"  Aula 05: {len(prs.slides)} slides -> aula-05-backtest-v1/")


# ══════════════════════════════════════════════════════════
# AULA 06 — ALOCACAO
# ══════════════════════════════════════════════════════════
def gerar_aula06():
    prs = new_prs()
    capa(prs, "06", "Alocacao: Como Pesar as Acoes do Portfolio",
         "O sinal diz QUAIS acoes comprar — a alocacao diz QUANTO de cada uma",
         ["Por que a escolha do esquema de pesos importa tanto quanto o sinal",
          "Equal-weight (1/N): simplicidade, robustez e o resultado surpresa de DeMiguel (2009)",
          "Vol-weight: reduzir contribuicao de ativos mais arriscados ao risco total",
          "Otimizacao de Markowitz: teoria da fronteira eficiente vs pratica",
          "Comparar as tres abordagens no mesmo sinal e tirar conclusoes"])

    agenda(prs,
        ["Equal-weight: por que 1/N supera MV fora da amostra",
         "Vol-weight: logica do risk parity simplificado",
         "Markowitz: onde a teoria encontra a realidade",
         "Turnover: custo oculto da rebalanceamento"],
        ["Implementar equal-weight (baseline)",
         "Implementar vol-weight (1/volatilidade normalizada)",
         "Implementar MV otimizado (scipy minimize)",
         "Comparar metricas das tres abordagens",
         "Analisar turnover e concentracao de cada esquema"])

    # Slide 3: Esquemas de pesos
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Esquemas de Alocacao — Tres Abordagens",
           "Da mais simples a mais sofisticada — nem sempre a mais complexa vence")
    txt(s, 0.5, 1.28, 12.3, 0.42,
        "A alocacao de capital entre os ativos selecionados e uma decisao separada do sinal. "
        "Diferentes esquemas de pesagem podem ter grande impacto no Sharpe, drawdown e turnover "
        "— e a melhor escolha frequentemente nao e a mais matematicamente sofisticada.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.83, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 1.98, 3.98, 4.82, "Equal-Weight (1/N)", GREEN,
          ["Peso de cada acao: 1 / n_acoes_selecionadas",
           "Maxima diversificacao entre os selecionados",
           "DeMiguel et al. (2009): supera MV em 14 datasets",
           "Robusto a erros de estimacao de parametros",
           "Baixo turnover por ser simetrico e simples",
           "Nossa baseline: o benchmark de comparacao"])
    card2(s, 4.67, 1.98, 3.98, 4.82, "Vol-Weight (Risk Parity simples)", BLUE,
          ["Peso proporcional a 1/volatilidade_historica",
           "Acoes mais volateis recebem peso menor",
           "Logica: equalizar contribuicao ao risco (nao ao capital)",
           "Vol estimada: desvio-padrao dos ultimos 63 dias uteis",
           "Pesos normalizados: somam 1 (long-only, sem alavancagem)",
           "Reduz dependencia de acoes muito volateis (ex: small caps)"])
    card2(s, 8.95, 1.98, 3.98, 4.82, "Markowitz (MV Otimizado)", RED,
          ["Maximizar Sharpe: max E[r]/σ sujeito a w >= 0, sum(w) = 1",
           "Requer estimacao de vetor de retornos esperados e matriz Σ",
           "Problema: parametros estimados com erro amostral grande",
           "Resultado: pesos concentrados e instáveis entre periodos",
           "Turnover muito alto: custo de transacao elimina o alpha teorico",
           "Na pratica: regularizacao (Black-Litterman, shrinkage) e necessaria"])
    ref(s, "DeMiguel, V., Garlappi, L. & Uppal, R. (2009). Optimal Versus Naive Diversification. "
           "Review of Financial Studies, 22(5), 1915–1953")

    # Slide 4: Turnover
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Turnover — O Custo Oculto da Rebalanceamento",
           "O quanto o portfolio muda a cada mes determina o custo real da estrategia")
    txt(s, 0.5, 1.28, 12.3, 0.45,
        "Turnover e a fracao do portfolio que e comprada ou vendida em cada rebalanceamento. "
        "Um turnover de 50% significa que metade do portfolio e renovada. "
        "Este numero, multiplicado pelo custo por transacao, determina o drag de custo mensal.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.87, 12.3, 0.03, fill=GOLD)
    rect(s, 0.4, 2.02, 6.1, 4.78, fill=LIGHT)
    txt(s, 0.6, 2.1, 5.8, 0.38, "Calculo do Turnover", size=14, bold=True, color=NAVY)
    txt(s, 0.6, 2.56, 5.8, 2.3,
        "# Turnover mensal\nturnover_t = (pesos_t - pesos_{t-1}).abs().sum() / 2\n\n"
        "# Turnover anualizado\nturnover_anual = turnover_mensal.mean() * 12\n\n"
        "# Custo mensal estimado (drag)\ncusto_mensal = turnover_t * custo_por_turno\nret_liquido = ret_bruto - custo_mensal",
        size=11, color=NAVY)
    buls(s, 0.6, 5.0, 5.8, 1.8,
         ["Equal-weight: turnover tipico 20-40% ao mes (entrada/saida de acoes)",
          "Vol-weight: turnover similar ao EW, pesos flutuam com vol",
          "Markowitz: turnover muito alto (60-100%) — inviavel sem restricoes"],
         size=12, color=NAVY, spacing=6)
    rect(s, 6.8, 2.02, 6.1, 4.78, fill=LIGHT)
    txt(s, 7.0, 2.1, 5.8, 0.38, "Impacto nos Retornos Liquidos", size=14, bold=True, color=NAVY)
    buls(s, 7.0, 2.56, 5.8, 4.0,
         ["Custo por turno tipico no Brasil: 0,3% a 0,5% (corretagem + B3)",
          "Turnover de 30% x custo 0,4% = drag de 0,12% ao mes = 1,4% a.a.",
          "Parece pouco, mas corroi 20-40% do alpha bruto tipico",
          "Break-even: qual custo maximo antes do alpha virar zero?",
          "Estrategias de baixo turnover sao muito mais robustas a custos",
          "Pratica: testar diferentes frequencias (mensal vs bimestral)"],
         size=12, color=NAVY, spacing=6)
    ref(s, "Keim, D.B. & Madhavan, A. (1997). Transactions Costs and Investment Style. "
           "J. of Financial Economics, 46(3), 265–292")

    construir(prs,
        ["pesos_equal_weight(sinal, n_top=10) — 1/N para top N acoes",
         "pesos_vol_weight(sinal, ret_diarios, n_top=10) — 1/vol, normalizado",
         "pesos_markowitz(sinal, ret_mensais, n_top=10) — scipy optimize + restricoes",
         "calcular_turnover(pesos_df) — |Δpesos| mensal e anualizado",
         "comparar_alocacoes(resultados_dict) — tabela e graficos lado a lado"],
        ["dados/retornos_diarios_limpo.parquet",
         "dados/retornos_mensais_limpo.parquet",
         "dados/sinal_v1.parquet"],
        ["dados/pesos_v1.parquet (atualizado com as 3 versoes)",
         "dados/retorno_carteira_v1.parquet (atualizado)"])

    out = os.path.join(BASE, "aula-06-alocacao", "slides-aula-06-alocacao.pptx")
    prs.save(out)
    print(f"  Aula 06: {len(prs.slides)} slides -> aula-06-alocacao/")


# ══════════════════════════════════════════════════════════
# AULA 07 — SINAL v2
# ══════════════════════════════════════════════════════════
def gerar_aula07():
    prs = new_prs()
    capa(prs, "07", "Sinal v2: Momentum Ajustado por Volatilidade",
         "Melhorando o sinal: retorno bruto ignora risco — vol-adjusted momentum corrige isso",
         ["Limitacao do sinal v1: acoes mais arriscadas tem retornos brutos maiores",
          "Volatility-adjusted momentum: dividir o sinal pelo risco historico",
          "Intuicao economica: comparar acoes em base de risco equivalente",
          "Comparacao rigorosa: IC, Sharpe e estabilidade v1 vs v2",
          "Escolha de parametros: janela de volatilidade e impacto no sinal"])

    agenda(prs,
        ["Por que retorno bruto e uma comparacao injusta entre acoes",
         "Vol-adjusted momentum: formula e intuicao",
         "Janela de volatilidade: qual o horizonte certo?",
         "Como comparar dois sinais de forma justa"],
        ["Calcular vol rolling de 63 dias uteis por acao",
         "Construir sinal_v2 = sinal_v1 / vol_rolling",
         "Comparar IC de v1 vs v2 ao longo do tempo",
         "Backtest completo do sinal v2",
         "Tabela comparativa de metricas v1 vs v2"])

    # Slide 3: Problema do sinal v1
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Limitacao do Sinal v1 — Retorno Bruto Ignora Risco",
           "Uma acao que subiu 30% com vol de 50% e diferente de uma que subiu 30% com vol de 15%")
    txt(s, 0.5, 1.28, 12.3, 0.45,
        "O sinal v1 ranqueia acoes pelo retorno absoluto dos ultimos 12-1 meses. "
        "Mas acoes mais arriscadas (maior volatilidade) tendem a ter retornos brutos maiores. "
        "Ao selecionar as top 10 por retorno bruto, podemos estar selecionando as mais arriscadas — "
        "e pagando por esse risco na forma de maior volatilidade do portfolio.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.87, 12.3, 0.03, fill=GOLD)
    rect(s, 0.4, 2.02, 5.9, 4.78, fill=LIGHT)
    txt(s, 0.6, 2.1, 5.6, 0.38, "Exemplo Numerico", size=14, bold=True, color=NAVY)
    buls(s, 0.6, 2.56, 5.7, 4.0,
         ["Acao A: retorno 12-1 = +40%, vol anual = 60%  →  sinal_v1 = 0,40",
          "Acao B: retorno 12-1 = +30%, vol anual = 20%  →  sinal_v1 = 0,30",
          "Ranking v1: A > B — A entra no portfolio",
          "Mas: A gerou 40%/60% = 0,67 Sharpe-like vs B com 30%/20% = 1,50",
          "Sinal v2 correto: A = 0,40/0,60 = 0,67; B = 0,30/0,20 = 1,50",
          "Ranking v2: B > A — B e mais atraente por unidade de risco"],
         size=12, color=NAVY, spacing=7)
    rect(s, 6.6, 2.02, 6.35, 4.78, fill=LIGHT)
    txt(s, 6.8, 2.1, 6.1, 0.38, "Sinal v2 — Formula e Implementacao", size=14, bold=True, color=NAVY)
    txt(s, 6.8, 2.56, 6.1, 2.0,
        "# Volatilidade rolling 63 dias uteis (anualizada)\nvol_rolling = (\n    ret_diarios\n    .rolling(63)\n    .std()\n    * np.sqrt(252)\n)\n\n"
        "# Resampled para frequencia mensal\nvol_mensal = vol_rolling.resample('ME').last()\n\n"
        "# Sinal vol-ajustado\nsinal_v2 = sinal_v1 / vol_mensal",
        size=11, color=NAVY)
    buls(s, 6.8, 4.65, 6.1, 2.0,
         ["63 dias uteis ~ 3 meses: janela tipica na literatura",
          "rolling na freq diaria: mais dados, estimativa mais estavel",
          ".last() no resample: pega vol mais recente de cada mes",
          "Normalizar apos calcular: pesos somam 1"],
         size=12, color=NAVY, spacing=6)
    ref(s, "Moskowitz, T., Ooi, Y.H. & Pedersen, L.H. (2012). Time Series Momentum. "
           "J. of Financial Economics, 104(2), 228–250")

    # Slide 4: Comparacao v1 vs v2
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Comparando v1 vs v2 — Como Fazer de Forma Justa",
           "Mesmas datas, mesmo benchmark, mesmo periodo — so o sinal muda")
    txt(s, 0.5, 1.28, 12.3, 0.42,
        "Para comparar dois sinais de forma justa, todos os outros elementos devem ser identicos: "
        "universo de acoes, esquema de alocacao, periodo de analise, benchmark, e metricas. "
        "A unica variavel e o sinal — assim qualquer diferenca de performance e atribuivel a ele.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.84, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 1.99, 3.98, 4.82, "O que Comparar", NAVY,
          ["IC medio e t-stat do IC: poder preditivo bruto",
           "Sharpe, Sortino, MDD, Calmar: qualidade do retorno",
           "Turnover: custo implicito de transacao",
           "Correlacao dos retornos: os sinais concordam?",
           "Estabilidade temporal: qual e mais consistente?"])
    card2(s, 4.67, 1.99, 3.98, 4.82, "Hipoteses do Comparativo", BLUE,
          ["Mesmo universo: top 10 acoes por sinal",
           "Mesma alocacao: equal-weight (isolar efeito do sinal)",
           "Mesmo periodo: datas com dados validos para ambos",
           "Mesmo benchmark: retorno do IBOVESPA (BOVA11 proxy)",
           "Mesmas metricas: funcao calcular_metricas() igual"])
    card2(s, 8.95, 1.99, 3.98, 4.82, "O que Esperar", GREEN,
          ["v2 tende a ter Sharpe ligeiramente superior ao v1",
           "v2 pode ter menor MDD (menos exposicao a acoes muito volateis)",
           "Turnover similar: mudanca e no sinal, nao na selecao de N",
           "Correlacao v1 vs v2 alta (0,7+): sinais concordam na maior parte",
           "Melhora pode ser pequena — o que importa e a direcao"])
    ref(s, "Israel, R., Moskowitz, T. (2013). The Role of Shorting, Firm Size, and Time on Market Anomalies. "
           "J. of Financial Economics, 108(2), 275–301")

    construir(prs,
        ["calcular_vol_rolling(ret_diarios, janela=63) — std anualizado por acao",
         "calcular_sinal_v2(sinal_v1, vol_mensal) — divisao elemento a elemento",
         "comparar_ic(sinal_v1, sinal_v2, ret_mensais) — IC serie temporal de cada",
         "backtest_sinal(sinal, ret_mensais, n_top=10) — funcao generica reutilizavel",
         "tabela_comparativa(resultados_v1, resultados_v2) — side-by-side formatado"],
        ["dados/retornos_diarios_limpo.parquet",
         "dados/retornos_mensais_limpo.parquet",
         "dados/sinal_v1.parquet"],
        ["dados/sinal_v2.parquet",
         "dados/pesos_sinal_v2.parquet",
         "dados/retorno_carteira_sinal_v2.parquet"])

    out = os.path.join(BASE, "aula-07-sinal-v2", "slides-aula-07-sinal-v2.pptx")
    prs.save(out)
    print(f"  Aula 07: {len(prs.slides)} slides -> aula-07-sinal-v2/")


# ══════════════════════════════════════════════════════════
# AULA 08 — BACKTEST RIGOROSO
# ══════════════════════════════════════════════════════════
def gerar_aula08():
    prs = new_prs()
    capa(prs, "08", "Backtest Rigoroso: 6 Armadilhas Fatais",
         "A diferenca entre um backtest que engana e um que convence o juri",
         ["Look-ahead bias: o erro que fabrica performance inexistente",
          "Overfitting: ajustar parametros ate o passado ficar perfeito",
          "Multiple testing: p-hacking e o Deflated Sharpe Ratio (Bailey & LdP, 2014)",
          "Custos de transacao: o assassino silencioso do alpha",
          "Walk-forward validation: a unica forma honesta de estimar performance OOS",
          "Como apresentar um backtest rigoros ao juri"])

    agenda(prs,
        ["Look-ahead bias: definicao e como contamina",
         "Overfitting: fitting noise vs signal",
         "Multiple testing e DSR: correcao estatistica",
         "Walk-forward: treino/teste temporal"],
        ["Demonstrar look-ahead bias (antes e depois)",
         "Sweep de parametros: 35 janelas, Sharpe por janela",
         "Simulacao de 1000 estrategias aleatorias (DSR)",
         "Modelar 4 cenarios de custo de transacao",
         "Walk-forward 48m treino / 12m teste rolling"])

    # Slide 3: Look-ahead bias
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Look-Ahead Bias — O Erro que Fabrica Alpha",
           "Usar informacao futura para tomar decisoes passadas: o mais perigoso dos erros")
    txt(s, 0.5, 1.28, 12.3, 0.45,
        "Look-ahead bias ocorre quando o sinal ou a decisao de portfolio usa informacao "
        "que nao estaria disponivel no momento da decisao. O resultado e um backtest "
        "artificialmente otimista que nao se replica na pratica.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.87, 12.3, 0.03, fill=GOLD)
    rect(s, 0.4, 2.02, 6.1, 4.78, fill=LIGHT)
    txt(s, 0.6, 2.1, 5.8, 0.38, "Como Acontece na Pratica", size=14, bold=True, color=NAVY)
    buls(s, 0.6, 2.56, 5.8, 4.2,
         ["ERRADO: sinal_t usa retorno DE t (inclui informacao de t)",
          "ERRADO: sinal = ret.rolling(12).sum()  # sem shift — contamina!",
          "CERTO: sinal_t usa retornos ate t-2 (shift garante isso)",
          "CERTO: sinal = ret.shift(2).rolling(11).sum()",
          "Outro exemplo: usar preco de fechamento de t para executar em t",
          "Correto: executar no fechamento de t+1 (next close)"],
         size=12, color=NAVY, spacing=7)
    rect(s, 6.8, 2.02, 6.1, 4.78, fill=LIGHT)
    txt(s, 7.0, 2.1, 5.8, 0.38, "Impacto Quantificado", size=14, bold=True, color=NAVY)
    buls(s, 7.0, 2.56, 5.8, 4.2,
         ["Backtest com look-ahead: Sharpe tipicamente 50-200% maior que o correto",
          "A diferenca e maior em mercados mais ineficientes (small caps, emergentes)",
          "Detectar: comparar performance IS com e sem o shift()",
          "Regra de ouro: trace cada dado do sinal ate sua fonte historica",
          "Em code review: procure .rolling() sem .shift() adequado antes",
          "Walk-forward nao corrige look-ahead — so o design correto do sinal corrige"],
         size=12, color=NAVY, spacing=7)
    ref(s, "Prado, M. Lopez de (2018). Advances in Financial Machine Learning, Cap. 11. Wiley")

    # Slide 4: Overfitting e DSR
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Overfitting e Deflated Sharpe Ratio",
           "Ajustar parametros ao passado fabrica performance — o DSR corrige o p-value")
    txt(s, 0.5, 1.28, 12.3, 0.42,
        "Com N parametros para escolher, e possivel encontrar uma combinacao que parece excelente "
        "historicamente por puro acaso (data snooping). O Deflated Sharpe Ratio de Bailey & "
        "Lopez de Prado (2014) calcula o Sharpe minimo necessario dado o numero de testes realizados.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.84, 12.3, 0.03, fill=GOLD)
    rect(s, 0.4, 1.99, 6.1, 4.82, fill=LIGHT)
    txt(s, 0.6, 2.07, 5.8, 0.38, "Deflated Sharpe Ratio (DSR)", size=14, bold=True, color=NAVY)
    txt(s, 0.6, 2.53, 5.8, 2.8,
        "# Sharpe minimo necessario dado N testes (Bailey & LdP 2014)\nfrom scipy.stats import norm\nimport numpy as np\n\nEULER = 0.5772156649\n\ndef sharpe_minimo(n_testes, sharpe_std=1.0):\n    sr = sharpe_std * (\n        (1 - EULER) * norm.ppf(1 - 1/n_testes)\n        + EULER * norm.ppf(1 - 1/(n_testes * np.e))\n    )\n    return sr\n\n# Exemplo: 35 testes -> SR minimo ~ 2.5",
        size=10, color=NAVY)
    buls(s, 0.6, 5.45, 5.8, 1.35,
         ["Quanto mais testes, maior o SR minimo necessario",
          "Fundamento: correcao de Bonferroni probabilistica"],
         size=12, color=NAVY, spacing=6)
    rect(s, 6.8, 1.99, 6.1, 4.82, fill=LIGHT)
    txt(s, 7.0, 2.07, 5.8, 0.38, "Como Detectar e Evitar", size=14, bold=True, color=NAVY)
    buls(s, 7.0, 2.53, 5.8, 4.2,
         ["Registrar TODOS os parametros testados (nao esconder os que falharam)",
          "Escolha de parametros deve ter fundamento economico ANTES do teste",
          "Exemplo: janela de 12 meses vem de Jegadeesh & Titman (1993) — nao de sweep",
          "Sweep de sensibilidade: mostrar que performance e robusta ao redor do parametro escolhido",
          "Simulacao de estrategias aleatorias: quantas batem nosso Sharpe por puro acaso?",
          "Walk-forward OOS: a validacao mais importante de todas"],
         size=12, color=NAVY, spacing=6)
    ref(s, "Bailey, D.H. & Lopez de Prado, M. (2014). The Deflated Sharpe Ratio. "
           "Financial Analysts Journal, 70(5), 94–107")

    # Slide 5: Walk-forward e custos
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Walk-Forward Validation e Custos de Transacao",
           "A validacao mais credivel e o drag de custos que o backtest simples ignora")
    rect(s, 0.5, 1.25, 12.3, 0.03, fill=GOLD)
    rect(s, 0.4, 1.35, 6.1, 5.45, fill=LIGHT)
    txt(s, 0.6, 1.42, 5.8, 0.38, "Walk-Forward Validation", size=14, bold=True, color=NAVY)
    buls(s, 0.6, 1.88, 5.8, 4.7,
         ["Divisao temporal: nao embaralhar dados (a diferenca do k-fold classico)",
          "Janela de treino: 48 meses (4 anos) — suficiente para estimar o sinal",
          "Janela de teste: 12 meses (1 ano) — OOS real, nunca visto",
          "Rolling: deslizar janela um mes de cada vez, concatenar resultados OOS",
          "Concatenar OOS: retornos de teste formam serie continua",
          "Comparar IS vs OOS: degradacao esperada ~30-50%; se maior, overfitting",
          "Regra: apresentar SO os resultados OOS ao juri — IS e auxiliar"],
         size=12, color=NAVY, spacing=6)
    rect(s, 6.8, 1.35, 6.1, 5.45, fill=LIGHT)
    txt(s, 7.0, 1.42, 5.8, 0.38, "Custos de Transacao", size=14, bold=True, color=NAVY)
    buls(s, 7.0, 1.88, 5.8, 4.7,
         ["Componentes: corretagem + taxa B3 + spread bid-ask + imposto",
          "Estimativa conservadora: 0,3% a 0,5% por turno (ida ou volta)",
          "Custo mensal = turnover_mensal x custo_por_turno",
          "Ret liquido = ret bruto - custo_mensal",
          "Testar 4 cenarios: 0%, 0,3%, 0,5%, 1,0% por turno",
          "Break-even: qual custo maximo que zera o alpha?",
          "Estrategia robusta: positiva ate em cenario de custo pessimista"],
         size=12, color=NAVY, spacing=6)
    ref(s, "Bailey & Lopez de Prado (2014); Keim & Madhavan (1997). J. of Financial Economics, 46(3), 265–292")

    construir(prs,
        ["demonstrar_lookahead(ret_mensais) — comparar sinal correto vs contaminado",
         "sweep_parametros(ret_mensais, janelas=range(2,37)) — Sharpe x janela",
         "simulacao_aleatorias(ret_mensais, n=1000) — distribuicao empirica de Sharpes",
         "sharpe_minimo_necessario(n_testes) — formula DSR de Bailey & LdP (2014)",
         "modelar_custos(ret_bruto, turnover, cenarios=[0, 0.003, 0.005, 0.01])",
         "walk_forward_momentum(ret_mensais, treino=48, teste=12) — retornos OOS"],
        ["dados/retornos_diarios_limpo.parquet",
         "dados/retornos_mensais_limpo.parquet",
         "dados/sinal_v2.parquet"],
        ["dados/retorno_walkforward_liquido.parquet"])

    out = os.path.join(BASE, "aula-08-backtest-rigoroso", "slides-aula-08-backtest-rigoroso.pptx")
    prs.save(out)
    print(f"  Aula 08: {len(prs.slides)} slides -> aula-08-backtest-rigoroso/")


# ══════════════════════════════════════════════════════════
# AULA 09 — GENAI + ANALISE
# ══════════════════════════════════════════════════════════
def gerar_aula09():
    prs = new_prs()
    capa(prs, "09", "GenAI: Analise Automatizada com Claude",
         "Usar LLMs para acelerar analise quantitativa — com limites claros",
         ["O que LLMs realmente fazem: predicao de tokens, nao 'inteligencia'",
          "API da Anthropic: autenticacao, chamadas, modelos disponiveis",
          "Prompt engineering: 4 principios para resultados confiaveis",
          "O que LLMs fazem bem (e mal) em analise quantitativa",
          "Workflow: metricas → prompt → rascunho → revisao humana"])

    agenda(prs,
        ["Como LLMs funcionam internamente (intuicao)",
         "API Anthropic: client, messages, system prompt",
         "4 principios de prompt engineering",
         "Limites dos LLMs em financas quantitativas"],
        ["Configurar API key e instanciar client",
         "Funcao chamar_claude() generica e reutilizavel",
         "Gerar analise de cada secao do relatorio",
         "Critica automatizada (juri virtual)",
         "Refinamento multi-turn: critica -> reescrita"])

    # Slide 3: Como LLMs funcionam
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Como LLMs Funcionam — Intuicao Essencial",
           "Entender o mecanismo evita tanto subestimar quanto superestimar a ferramenta")
    txt(s, 0.5, 1.28, 12.3, 0.42,
        "Large Language Models sao redes neurais treinadas para prever o proximo token "
        "em uma sequencia de texto. Essa tarefa aparentemente simples, em escala massiva, "
        "emerge em capacidades surpreendentes — mas com limitacoes estruturais importantes.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.84, 12.3, 0.03, fill=GOLD)
    card2(s, 0.4, 1.99, 3.98, 4.82, "O que LLMs fazem bem", GREEN,
          ["Sintetizar e estruturar texto com qualidade",
           "Traduzir metricas numericas em linguagem natural",
           "Seguir formatos especificos (JSON, Markdown, tabelas)",
           "Gerar rascunhos rapidos de secoes de relatorio",
           "Identificar pontos fracos em argumentacoes",
           "Simular perguntas de juri e avaliar respostas"])
    card2(s, 4.67, 1.99, 3.98, 4.82, "O que LLMs NAO fazem", RED,
          ["Calcular Sharpe ratios ou qualquer metrica numerica (alucinam!)",
           "Acessar dados de mercado em tempo real",
           "Garantir que fatos citados sao corretos (hallucination)",
           "Substituir analise estatistica rigorosa",
           "Prever precos ou retornos futuros",
           "Validar se o backtest esta correto logicamente"])
    card2(s, 8.95, 1.99, 3.98, 4.82, "Regra de Ouro", NAVY,
          ["Sempre passe as metricas calculadas no Python como input",
           "Nunca peca para o LLM calcular — so para interpretar",
           "Valide toda afirmacao factual gerada pelo modelo",
           "Use o LLM para rascunho e aceleracao, nao como oraculo",
           "O relatorio final deve refletir SUA analise, nao so do LLM",
           "Citacoes: verifique manualmente — LLMs inventam referencias"])
    ref(s, "Vaswani, A. et al. (2017). Attention Is All You Need. NeurIPS; "
           "Anthropic (2024). Claude Model Card. anthropic.com")

    # Slide 4: Prompt Engineering
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Prompt Engineering — 4 Principios Fundamentais",
           "A qualidade do output depende diretamente da qualidade do input")
    txt(s, 0.5, 1.28, 12.3, 0.42,
        "Prompt engineering e a pratica de estruturar instrucoes para LLMs de forma a obter "
        "outputs uteis, precisos e no formato desejado. Nao e magia — e comunicacao clara "
        "com um sistema que nao tem contexto alem do que voce fornece.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.84, 12.3, 0.03, fill=GOLD)
    principios = [
        ("1. Contexto Rico", NAVY,
         "Quanto mais contexto relevante, melhor o output.\n"
         "Incluir: objetivo da analise, publico-alvo, metricas calculadas, restricoes.\n"
         "Exemplo: 'Voce e um analista quant senior. A estrategia e momentum 12-1 "
         "no IBOVESPA com Sharpe de 0,85 e MDD de -18%. O leitor e o juri do Desafio Itau.'"),
        ("2. Formato Especificado", BLUE,
         "Especificar o formato de saida desejado reduz pos-processamento.\n"
         "Exemplos: 'Responda em JSON com campos {secao, conteudo, pontos_fracos}'\n"
         "ou 'Use markdown com cabecalhos H2, bullets e no maximo 300 palavras'.\n"
         "LLMs sao bons em seguir formatos quando explicitamente pedido."),
        ("3. Dados no Prompt", GREEN,
         "Nunca peca para o LLM calcular — passe os numeros calculados.\n"
         "Errado: 'Analise a performance da estrategia'\n"
         "Certo: 'Sharpe=0.85, Sortino=1.12, MDD=-18%, CAGR=14.2%, Calmar=0.79. "
         "Benchmark IBOV: Sharpe=0.41, MDD=-35%. Interprete esses resultados.'"),
        ("4. Chain of Thought", ORAN,
         "Para analises complexas, pedir raciocinio passo a passo melhora a qualidade.\n"
         "'Antes de concluir, liste os pontos fortes e fracos da estrategia, "
         "depois avalie a consistencia dos resultados, e finalmente escreva a conclusao.'\n"
         "Refinamento multi-turn: gerar → criticar → reescrever."),
    ]
    for i, (title, color, body) in enumerate(principios):
        row, col = i // 2, i % 2
        x = 0.4 + col * 6.45
        y = 2.02 + row * 2.52
        rect(s, x, y, 6.2, 2.35, fill=LIGHT)
        rect(s, x, y, 6.2, 0.4, fill=color)
        txt(s, x + 0.15, y + 0.06, 5.9, 0.32, title, size=13, bold=True, color=WHITE)
        txt(s, x + 0.15, y + 0.5, 5.9, 1.72, body, size=11, color=NAVY)
    ref(s, "Brown, T. et al. (2020). Language Models are Few-Shot Learners. NeurIPS; "
           "Wei, J. et al. (2022). Chain-of-Thought Prompting Elicits Reasoning in LLMs. NeurIPS")

    construir(prs,
        ["chamar_claude(prompt, system, modelo, max_tokens) — wrapper da API generalizavel",
         "metricas_para_texto(metricas_dict) — serializar numeros para o prompt",
         "gerar_secao(secao, metricas, system_prompt) — gerar cada parte do relatorio",
         "criticar_secao(texto, system_juri) — simular juri avaliando a secao",
         "conversa_refinamento(secao, critica) — multi-turn: gerar → criticar → reescrever",
         "montar_relatorio_completo() — juntar todas as secoes em Markdown"],
        ["dados/retorno_carteira_sinal_v2.parquet",
         "dados/retorno_walkforward_liquido.parquet"],
        ["relatorio_draft.md (arquivo texto com rascunho completo)"])

    out = os.path.join(BASE, "aula-09-genai-analise", "slides-aula-09-genai-analise.pptx")
    prs.save(out)
    print(f"  Aula 09: {len(prs.slides)} slides -> aula-09-genai-analise/")


# ══════════════════════════════════════════════════════════
# AULA 10 — RELATORIO + DEFESA
# ══════════════════════════════════════════════════════════
def gerar_aula10():
    prs = new_prs()
    capa(prs, "10", "Relatorio Final e Defesa Oral",
         "O produto final do Intensivao: apresentar e defender a estrategia ao juri",
         ["Estrutura do relatorio tecnico quant: o que vai em cada secao",
          "Tear sheet: o cartao de visita visual da estrategia (5 paineis)",
          "Defesa oral: como responder perguntas do juri com fundamento",
          "Rubrica de avaliacao: 100 pontos distribuidos em 7 criterios",
          "Cross-review: avaliar os outros times para afiar o proprio argumento"])

    agenda(prs,
        ["Estrutura do relatorio quant profissional",
         "O que o juri realmente avalia (e como pontuar alto)",
         "Como defender escolhas metodologicas",
         "Erros fatais na defesa oral"],
        ["Tear sheet final: 5 paineis em uma figura",
         "Mapeamento de criterios de avaliacao por aula",
         "Rubrica de cross-review (100 pontos)",
         "Preparar 5 respostas para perguntas dificeis",
         "Checklist final de entrega"])

    # Slide 3: Estrutura do relatorio
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Estrutura do Relatorio Tecnico Quant",
           "Um relatorio profissional tem secoes definidas — cada uma responde a uma pergunta do juri")
    txt(s, 0.5, 1.28, 12.3, 0.42,
        "Um relatorio quant bem estruturado nao apenas apresenta resultados — ele conta uma historia: "
        "'Identificamos uma anomalia empiricamente suportada, implementamos de forma rigoros e "
        "obtemos resultado robusto fora da amostra, defensavel diante de custos reais.'",
        size=13, color=NAVY)
    rect(s, 0.5, 1.84, 12.3, 0.03, fill=GOLD)
    secoes = [
        ("1. Resumo Executivo", GOLD,
         "1 pagina maxima. A resposta para: qual e a estrategia, qual o resultado principal, "
         "por que e confiavel? Escreva por ultimo, mas posicione primeiro."),
        ("2. Hipotese e Fundamentacao", NAVY,
         "Por que momentum? Qual a teoria economica? Quais as referencias academicas primarias? "
         "Mostrar que a tese tem base — nao e mineracao de dados."),
        ("3. Dados e Metodologia", BLUE,
         "Universo, periodo, fonte, limpeza realizada, sinal exato, esquema de alocacao. "
         "Detalhe suficiente para que qualquer pessoa replique o backtest."),
        ("4. Resultados do Backtest", GREEN,
         "Metricas completas (IS e OOS separados). Comparacao com benchmark. "
         "Analise de custos. Walk-forward. Graficos de retorno acumulado e drawdown."),
        ("5. Analise de Risco", RED,
         "Sensibilidade a parametros. Cenarios de stress. Periodos de underperformance. "
         "Limitacoes explicitas do modelo. O que pode dar errado na pratica."),
        ("6. Conclusao", PURP,
         "Sintetizar os achados. Contribuicao da estrategia. "
         "Proximos passos se fosse implementar em producao."),
    ]
    for i, (title, color, body) in enumerate(secoes):
        row, col = i // 3, i % 3
        x = 0.4 + col * 4.25
        y = 2.02 + row * 2.42
        rect(s, x, y, 4.05, 2.25, fill=LIGHT)
        rect(s, x, y, 4.05, 0.4, fill=color)
        txt(s, x + 0.12, y + 0.06, 3.82, 0.32, title, size=11, bold=True, color=WHITE)
        txt(s, x + 0.12, y + 0.5, 3.82, 1.6, body, size=10.5, color=NAVY)
    ref(s, "Lo, A.W. (2016). What is an index? J. of Portfolio Management, 42(2), 21–36; "
           "Grinold, R. & Kahn, R. (1999). Active Portfolio Management. McGraw-Hill")

    # Slide 4: Defesa oral
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Defesa Oral — Como Responder Perguntas do Juri",
           "Cada pergunta do juri tem uma estrutura de resposta ideal: fundamento → alternativa → sensibilidade")
    txt(s, 0.5, 1.28, 12.3, 0.42,
        "O juri vai desafiar cada escolha metodologica. A resposta ideal tem tres partes: "
        "(1) o fundamento academico ou economico da escolha, "
        "(2) as alternativas que foram consideradas e por que foram descartadas, "
        "(3) a analise de sensibilidade mostrando que o resultado e robusto.",
        size=13, color=NAVY)
    rect(s, 0.5, 1.84, 12.3, 0.03, fill=GOLD)
    perguntas = [
        ("Por que janela 12-1?",
         "Fundamento: Jegadeesh & Titman (1993) e o padrao da literatura — nao escolhemos para maximizar Sharpe.\n"
         "Alternativas: testamos 3-1 a 36-1; performance e estavel para J entre 9 e 15 meses.\n"
         "Sensibilidade: grafico de Sharpe x janela mostra platô robusto ao redor de 12."),
        ("Por que equal-weight?",
         "Fundamento: DeMiguel et al. (2009) mostram que 1/N supera MV em 14 datasets fora da amostra.\n"
         "Alternativas: testamos vol-weight e MV; vol-weight levemente melhor, MV pior apos custos.\n"
         "Sensibilidade: diferenca de Sharpe entre EW e VW e pequena — escolhemos robustez."),
        ("O backtest tem look-ahead?",
         "Fundamento: sinal usa ret.shift(2).rolling(11) — dados disponiveis ate t-2 no maximo.\n"
         "Alternativas: testamos sem shift — contamina e Sharpe sobe 40%; descartamos.\n"
         "Sensibilidade: walk-forward OOS confirma performance sem acesso a dados futuros."),
        ("Como sabem que nao e overfitting?",
         "Fundamento: janela escolhida por teoria, nao por sweep; DSR calculado com N=35 testes.\n"
         "Alternativas: simulamos 1000 estrategias aleatorias — nosso Sharpe supera 95% delas.\n"
         "Sensibilidade: performance OOS (walk-forward) e 70% da IS — degradacao esperada."),
        ("A estrategia funciona com custos reais?",
         "Fundamento: modelamos 4 cenarios: 0%, 0,3%, 0,5%, 1,0% por turno.\n"
         "Alternativas: break-even de custo e X% — acima disso alpha desaparece.\n"
         "Sensibilidade: com custo de 0,3% (conservador), Sharpe liquido ainda e positivo."),
    ]
    for i, (q, a) in enumerate(perguntas):
        y = 2.05 + i * 0.98
        rect(s, 0.4, y, 12.5, 0.9, fill=LIGHT if i % 2 == 0 else WHITE)
        txt(s, 0.6, y + 0.05, 3.5, 0.35, q, size=11, bold=True, color=NAVY)
        txt(s, 4.3, y + 0.05, 8.5, 0.78, a, size=10, color=NAVY)

    # Slide 5: Checklist final
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    header(s, "Checklist Final de Entrega",
           "Tudo que o relatorio e a defesa precisam ter — marque antes de submeter")
    rect(s, 0.5, 1.25, 12.3, 0.03, fill=GOLD)
    itens_l = [
        "[ ] Hipotese economica clara com referencias academicas primarias",
        "[ ] Sinal implementado com shift() correto (sem look-ahead bias)",
        "[ ] Backtest com custos reais modelados (minimo 2 cenarios)",
        "[ ] Walk-forward OOS reportado separadamente do IS",
        "[ ] Deflated Sharpe Ratio calculado e reportado",
        "[ ] Analise de sensibilidade a parametros (janela, N acoes)",
        "[ ] Tear sheet com retorno acumulado e drawdown como figura",
    ]
    itens_r = [
        "[ ] Comparacao com benchmark IBOVESPA em todas as metricas",
        "[ ] Metricas: Sharpe, Sortino, MDD, Calmar, alpha, beta",
        "[ ] Periodos de underperformance discutidos honestamente",
        "[ ] Limitacoes do modelo explicitadas na conclusao",
        "[ ] Relatorio com revisao humana (nao so GenAI)",
        "[ ] 5 perguntas dificeis preparadas com estrutura fund→alt→sens",
        "[ ] Notebook roda do inicio ao fim sem erros antes de entregar",
    ]
    for i, item in enumerate(itens_l):
        c = GREEN if "[ ]" in item else NAVY
        txt(s, 0.5, 1.42 + i * 0.75, 6.3, 0.65, item, size=12, color=NAVY)
    for i, item in enumerate(itens_r):
        txt(s, 6.9, 1.42 + i * 0.75, 6.3, 0.65, item, size=12, color=NAVY)
    rect(s, 0.4, 6.75, 12.5, 0.55, fill=NAVY)
    txt(s, 0.6, 6.82, 12.0, 0.38,
        "Lembre: o juri avalia rigor metodologico e capacidade de defesa — nao so o resultado numerico.",
        size=13, bold=True, color=GOLD)

    construir(prs,
        ["tearsheet_final(retornos, benchmark, metricas) — figura 5-paineis completa",
         "plot_mapeamento_criterios() — cobertura de cada aula nos 7 criterios do juri",
         "rubrica_crossreview() — tabela 100 pontos para avaliar outras equipes",
         "preparar_defesa(perguntas, respostas) — formatar Q&A com estrutura fund→alt→sens",
         "checklist_final() — verificar todos os requisitos do relatorio"],
        ["dados/retorno_carteira_sinal_v2.parquet",
         "dados/retorno_walkforward_liquido.parquet",
         "relatorio_draft.md"],
        ["relatorio_final.md (versao revisada e completa)",
         "tearsheet_final.png",
         "crossreview_rubrica.png"])

    out = os.path.join(BASE, "aula-10-relatorio-defesa", "slides-aula-10-relatorio-defesa.pptx")
    prs.save(out)
    print(f"  Aula 10: {len(prs.slides)} slides -> aula-10-relatorio-defesa/")


# ══════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("Gerando slides das aulas 02 a 10...")
    gerar_aula02()
    gerar_aula03()
    gerar_aula04()
    gerar_aula05()
    gerar_aula06()
    gerar_aula07()
    gerar_aula08()
    gerar_aula09()
    gerar_aula10()
    print("Concluido! 9 arquivos .pptx gerados.")
