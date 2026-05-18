"""
Gerador de slides — Intensivão Quant AI | Aula 1: Kickoff
Execução: python gerar_slides_aula01.py
Saída:    slides-aula-01-kickoff.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta de cores ──────────────────────────────────────────
NAVY  = RGBColor(0x0D, 0x1B, 0x3E)
GOLD  = RGBColor(0xF5, 0xA6, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
GRAY  = RGBColor(0x77, 0x88, 0x9A)
GREEN = RGBColor(0x1E, 0x8B, 0x4C)
BLUE  = RGBColor(0x1A, 0x6E, 0xAE)

W, H = 13.33, 7.5
prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ── Primitivos ───────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp


def txt(slide, l, t, w, h, content, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, l, t, w, h, items, size=13, color=WHITE, spacing=8):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def ref(slide, citation):
    txt(slide, 0.3, 7.05, W - 0.5, 0.38,
        f"Ref.: {citation}", size=9, color=GRAY, italic=True)


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, 1.2, fill=NAVY)
    rect(slide, 0, 0, 0.22, 1.2, fill=GOLD)
    txt(slide, 0.42, 0.12, W - 0.6, 0.58, title,
        size=23, bold=True, color=WHITE)
    if subtitle:
        txt(slide, 0.42, 0.72, W - 0.6, 0.40, subtitle,
            size=12, color=LIGHT, italic=True)


def card(slide, l, t, w, h, title, body_items, bg=LIGHT, title_bg=NAVY, size=12):
    rect(slide, l, t, w, h, fill=bg)
    rect(slide, l, t, w, 0.42, fill=title_bg)
    txt(slide, l + 0.12, t + 0.07, w - 0.2, 0.34,
        title, size=12, bold=True, color=WHITE)
    bullets(slide, l + 0.12, t + 0.52, w - 0.2, h - 0.62,
            body_items, size=size, color=NAVY, spacing=5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, 0, 0.45, H, fill=GOLD)
rect(s, 0.45, 3.75, W - 0.45, 0.04, fill=GOLD)

txt(s, 0.85, 0.5, 11.5, 0.6,
    "IMPACTUFSC AR  —  QUANT DIRECTORATE", size=11, color=GOLD, italic=True)
txt(s, 0.85, 1.15, 11.5, 1.4,
    "Intensivão Quant AI", size=46, bold=True, color=WHITE)
txt(s, 0.85, 2.65, 11.5, 0.65,
    "Aula 1 — Kickoff", size=28, color=GOLD)
txt(s, 0.85, 3.35, 11.0, 0.38,
    "Fundamentos do Mercado Financeiro & Introdução à Estratégia Quantitativa",
    size=15, color=LIGHT)
txt(s, 0.85, 4.0, 9.0, 0.35,
    "Desafio Quant AI — Itaú Asset Management  |  2025", size=13, color=GRAY)

bullets(s, 0.85, 4.8, 11.0, 2.3,
        ["Universo: ~77 ações constituintes do IBOVESPA (dados via yfinance)",
         "Estratégia: momentum cross-sectional 12-1 meses (Jegadeesh & Titman, 1993)",
         "Objetivo: construir, backtester e defender uma estratégia quant completa em 10 aulas"],
        size=13, color=LIGHT, spacing=6)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Agenda — Aula 1", "O que vamos cobrir hoje")

items = [
    "1   O ecossistema financeiro — classes de ativos e seus papéis no portfólio",
    "2   Renda Fixa — yield, duration e risco de crédito",
    "3   Renda Variável — equity risk premium e por que ações sobem no longo prazo",
    "4   A Hipótese dos Mercados Eficientes e suas anomalias empíricas",
    "5   Factor Investing — do stock-picking a fatores sistemáticos (Fama-French, Carhart)",
    "6   O Efeito Momentum — fundamento acadêmico da nossa estratégia",
    "7   Behavioral Finance — por que o momentum persiste se os mercados são eficientes?",
    "8   O Desafio Itaú Asset Management — critérios e roadmap do Intensivão",
]
for i, item in enumerate(items):
    y = 1.38 + i * 0.68
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, 0.4, y, 12.5, 0.63, fill=bg)
    txt(s, 0.65, y + 0.12, 12.0, 0.42, item, size=15, color=NAVY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — O ECOSSISTEMA FINANCEIRO GLOBAL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "O Ecossistema Financeiro Global",
       "Onde o capital é alocado e qual a função de cada classe de ativo")

txt(s, 0.5, 1.3, 12.3, 0.5,
    "O sistema financeiro canaliza poupança para investimentos produtivos. Cada classe de ativo "
    "representa um contrato diferente de risco e retorno entre quem poupa e quem precisa de capital.",
    size=13, color=NAVY)

classes = [
    ("Renda Fixa", BLUE,
     ["~US$ 130 tri globalmente", "Governo e empresas tomam empréstimos",
      "Retorno contratual (cupom + principal)", "Risco: crédito, juros, inflação"]),
    ("Renda Variável", GREEN,
     ["~US$ 110 tri em market cap global", "Participação no capital de empresas",
      "Equity risk premium histórico ~5% a.a.", "Risco: mercado, empresa, setor"]),
    ("Derivativos", RGBColor(0x6C, 0x35, 0x83),
     ["Nocional >US$ 600 tri (BIS, 2024)", "Futuros, opções, swaps, forwards",
      "Hedge ou alavancagem de posições", "Risco: contraparte, liquidez, base"]),
    ("Alternativos", RGBColor(0xBF, 0x6A, 0x02),
     ["FIIs, PE, VC, Hedge Funds, Commodities", "Baixa correlação com RV e RF",
      "Illiquidity premium como compensação", "Risco: liquidez, valuation, lock-up"]),
]
for i, (title, color, body) in enumerate(classes):
    card(s, 0.35 + i * 3.15, 2.0, 3.0, 4.85, title, body,
         title_bg=color, size=13)

ref(s, "BIS Quarterly Review (2024); World Bank Financial Development Database (2024); "
       "Bodie, Kane & Marcus (2023). Investments, 13ª ed., McGraw-Hill")


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — RENDA FIXA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Renda Fixa — Fundamentos",
       "O mercado que empresta dinheiro ao governo e às empresas")

txt(s, 0.5, 1.3, 12.3, 0.55,
    "Ao comprar um título de renda fixa, o investidor empresta capital ao emissor em troca de uma "
    "remuneração pré-acordada (cupom) e da devolução do principal no vencimento. O retorno é "
    "'fixo' em estrutura, mas o preço de mercado do título flutua com as taxas de juros.",
    size=13, color=NAVY)

rect(s, 0.5, 2.0, 12.3, 0.03, fill=GOLD)

cols_rf = [
    ("Tipos no Brasil", BLUE,
     ["Tesouro Direto: LTN (pré-fixado), NTN-B (IPCA+), LFT (Selic)",
      "CDB e LCI/LCA: emissor bancário, garantia FGC até R$ 250k",
      "Debêntures: empresas não-financeiras, isenção IR (incentivadas)",
      "CRIs e CRAs: crédito imobiliário e agro, isentos para PF"]),
    ("Métricas Essenciais", NAVY,
     ["Yield to maturity (YTM): retorno total se mantido até vencimento",
      "Duration (Macaulay, 1938): prazo médio ponderado pelos fluxos",
      "Sensibilidade: ΔPreço ≈ −Duration × ΔJuros (modified duration)",
      "Spread de crédito: prêmio exigido acima do risco soberano"]),
    ("Risco × Retorno", RGBColor(0x1A, 0x6E, 0x4A),
     ["Tesouro Selic: risco mínimo, retorno CDI (~10,5% a.a. em 2024)",
      "NTN-B longa: proteção inflacionária, alta duration → vol de preço",
      "Debênture high yield: spread 3-8% sobre CDI, risco default real",
      "Regra geral: duration longa + rating baixo = maior volatilidade"]),
]
for i, (title, color, body) in enumerate(cols_rf):
    card(s, 0.4 + i * 4.18, 2.1, 3.98, 4.65, title, body,
         title_bg=color, size=12)

ref(s, "Macaulay, F.R. (1938). Some Theoretical Problems Suggested by the Movements of Interest Rates. NBER; "
       "Fabozzi, F.J. (2012). The Handbook of Fixed Income Securities, 8ª ed., McGraw-Hill")


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — RENDA VARIÁVEL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Renda Variável — Fundamentos",
       "Por que ações geram retorno superior no longo prazo? O equity risk premium")

txt(s, 0.5, 1.3, 12.3, 0.55,
    "Uma ação representa fração da propriedade de uma empresa. O acionista tem direito residual "
    "sobre os lucros (dividendos) e o patrimônio — mas arca com toda a variabilidade dos resultados. "
    "O equity risk premium (ERP) é a compensação exigida por esse risco adicional em relação ao ativo sem risco.",
    size=13, color=NAVY)

rect(s, 0.5, 2.0, 12.3, 0.03, fill=GOLD)

txt(s, 0.5, 2.1, 12.3, 0.38,
    "Equity Risk Premium histórico (1900–2023) — retorno real de ações acima dos títulos soberanos:",
    size=13, bold=True, color=NAVY)

erp = [("EUA", "5,5% a.a."), ("Brasil", "~6,0% a.a."),
       ("Reino Unido", "4,8% a.a."), ("Média Global", "~5,0% a.a.")]
for i, (pais, val) in enumerate(erp):
    col, row = i % 2, i // 2
    rect(s, 0.5 + col * 6.2, 2.58 + row * 1.45, 5.95, 1.3, fill=LIGHT)
    txt(s, 0.7 + col * 6.2, 2.65 + row * 1.45, 5.5, 0.38,
        pais, size=14, bold=True, color=NAVY)
    txt(s, 0.7 + col * 6.2, 3.0 + row * 1.45, 5.5, 0.45,
        val, size=22, bold=True, color=GREEN)

txt(s, 0.5, 5.55, 12.3, 0.38,
    "Veículos de investimento: ações individuais (B3/BVMF), ETFs (BOVA11, IVVB11), BDRs, fundos de ações.",
    size=13, color=NAVY)
txt(s, 0.5, 5.98, 12.3, 0.38,
    "Riscos: mercado (beta — risco sistemático não diversificável) + empresa (risco idiossincrático, diversificável).",
    size=13, color=NAVY)

ref(s, "Mehra, R. & Prescott, E.C. (1985). The Equity Premium: A Puzzle. J. of Monetary Economics, 15(2), 145–161; "
       "Dimson, Marsh & Staunton (2002). Triumph of the Optimists. Princeton University Press")


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — DERIVATIVOS E ALTERNATIVOS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Derivativos & Ativos Alternativos",
       "Instrumentos para gestão de risco e diversificação além do mercado tradicional")

# Derivativos — coluna esquerda
rect(s, 0.4, 1.3, 6.0, 5.95, fill=LIGHT)
rect(s, 0.4, 1.3, 6.0, 0.42, fill=RGBColor(0x6C, 0x35, 0x83))
txt(s, 0.6, 1.36, 5.6, 0.35, "Derivativos", size=15, bold=True, color=WHITE)

derivs = [
    ("Futuros",
     "Obrigação contratual de comprar ou vender um ativo-objeto em data e preço futuros "
     "definidos hoje. Principal uso: hedge de câmbio (exportadores), juros (Duration hedge) "
     "e commodities (agro). Negociados em bolsa (B3), com ajuste diário de margem."),
    ("Opções",
     "Direito (não obrigação) de comprar (call) ou vender (put) um ativo a preço fixado (strike) "
     "até ou na data de vencimento. Custo: prêmio pago na entrada. Apreçamento: modelo de "
     "Black & Scholes (1973) e extensões (árvores binomiais, Monte Carlo)."),
    ("Swaps",
     "Acordo de troca periódica de fluxos financeiros entre duas partes. Exemplo clássico: "
     "CDI × taxa pré-fixada. Empresas com passivo indexado ao IPCA usam swap para converter "
     "para CDI. Negociados OTC via B3/CETIP com câmaras de compensação."),
]
for i, (name, body) in enumerate(derivs):
    txt(s, 0.6, 1.85 + i * 1.58, 5.7, 0.3,
        name, size=13, bold=True, color=RGBColor(0x6C, 0x35, 0x83))
    txt(s, 0.6, 2.17 + i * 1.58, 5.7, 0.95,
        body, size=11, color=NAVY)

# Alternativos — coluna direita
rect(s, 6.93, 1.3, 6.0, 5.95, fill=LIGHT)
rect(s, 6.93, 1.3, 6.0, 0.42, fill=RGBColor(0xBF, 0x6A, 0x02))
txt(s, 7.13, 1.36, 5.6, 0.35, "Ativos Alternativos", size=15, bold=True, color=WHITE)

alts = [
    ("FIIs — Fundos de Investimento Imobiliário",
     "Cotas de portfólios de imóveis (lajes corporativas, shoppings, galpões logísticos). "
     "Distribuem rendimentos mensais (isentos de IR para PF). Negociados na B3 como ações. "
     "Métricas: dividend yield, cap rate, vacância física e financeira."),
    ("Private Equity e Venture Capital",
     "Participação em empresas de capital fechado. PE: empresas maduras em reestruturação. "
     "VC: startups em estágio inicial com potencial de crescimento exponencial. "
     "Prêmio de iliquidez: compensação por não poder sair a qualquer momento."),
    ("Commodities e Ouro",
     "Ativos reais: petróleo, minério de ferro, soja, milho, ouro. Hedge natural contra "
     "inflação. Ouro: store of value em crises (correlação negativa com risco). "
     "Acesso via ETFs (OZ1D11), fundos ou contratos futuros na B3."),
]
for i, (name, body) in enumerate(alts):
    txt(s, 7.13, 1.85 + i * 1.58, 5.7, 0.3,
        name, size=13, bold=True, color=RGBColor(0xBF, 0x6A, 0x02))
    txt(s, 7.13, 2.17 + i * 1.58, 5.7, 0.95,
        body, size=11, color=NAVY)

ref(s, "Black, F. & Scholes, M. (1973). The Pricing of Options and Corporate Liabilities. "
       "J. of Political Economy, 81(3), 637–654; ANBIMA (2024). Panorama de Mercado de Capitais")


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — TEORIA MODERNA DE PORTFÓLIO
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Teoria Moderna de Portfólio (MPT)",
       "A formalização matemática de 'não coloque todos os ovos na mesma cesta'")

txt(s, 0.5, 1.3, 12.3, 0.55,
    "Harry Markowitz (1952) revolucionou a gestão de investimentos ao mostrar que o que importa "
    "não é o risco individual de cada ativo, mas a contribuição ao risco total do portfólio — "
    "determinada pelas correlações. Diversificação reduz o risco sem necessariamente reduzir o retorno.",
    size=13, color=NAVY)

rect(s, 0.5, 2.0, 12.3, 0.03, fill=GOLD)

cols_mpt = [
    ("Conceito Central", NAVY,
     ["Retorno esperado: E[Rp] = Σ wᵢ · E[Rᵢ]",
      "Variância: σ²p = Σᵢ Σⱼ wᵢ wⱼ σᵢⱼ (covariâncias importam!)",
      "Covariância negativa → diversificação elimina risco idiossincrático",
      "Fronteira eficiente: máximo retorno para cada nível de risco",
      "Portfólio de mínima variância: mínimo risco independente de retorno"]),
    ("CAPM — Extensão do MPT", BLUE,
     ["Sharpe (1964): em equilíbrio geral, todos têm o portfólio de mercado",
      "E[Ri] = Rf + βi · (E[Rm] − Rf)  — equação do CAPM",
      "Beta: medida de risco sistemático (não diversificável)",
      "α (alpha) = retorno acima do previsto pelo CAPM dado o beta",
      "Fundamento do custo de capital (WACC) em valuation"]),
    ("Limitações Práticas", RGBColor(0xB0, 0x33, 0x2A),
     ["Parâmetros estimados com erro amostral significativo",
      "Portfólio ótimo instável: pequena mudança de input → grande rotação",
      "Retornos não gaussianos: fat tails subestimam risco de cauda",
      "DeMiguel et al. (2009): equal-weight (1/N) supera MV fora da amostra",
      "Na prática: restrições e regularização são necessárias"]),
]
for i, (title, color, body) in enumerate(cols_mpt):
    card(s, 0.4 + i * 4.18, 2.1, 3.98, 4.65, title, body,
         title_bg=color, size=12)

ref(s, "Markowitz, H. (1952). Portfolio Selection. J. of Finance, 7(1), 77–91; "
       "Sharpe, W.F. (1964). Capital Asset Prices. J. of Finance, 19(3), 425–442; "
       "DeMiguel, Garlappi & Uppal (2009). Rev. of Financial Studies, 22(5), 1915–1953")


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — HIPÓTESE DOS MERCADOS EFICIENTES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "A Hipótese dos Mercados Eficientes (HME)",
       "Se preços refletem toda informação disponível, estratégias ativas não geram alfa — mas será verdade?")

txt(s, 0.5, 1.3, 12.3, 0.5,
    "Eugene Fama (1970) definiu que um mercado é eficiente quando os preços dos ativos incorporam "
    "instantaneamente toda a informação relevante, tornando impossível obter retorno acima do risco "
    "de forma consistente. Fama propôs três formas de eficiência, com implicações distintas.",
    size=13, color=NAVY)

rect(s, 0.5, 1.95, 12.3, 0.03, fill=GOLD)

forms = [
    ("Forma Fraca", BLUE,
     "Preços incorporam todo o histórico de preços e volumes negociados.\n"
     "\nImplicação direta: análise técnica (chartismo, suporte/resistência) "
     "não gera alfa sistemático — toda informação passada já está no preço.\n"
     "\nEvidência: autocorrelação de curto prazo existe (Fama, 1988), "
     "mas é pequena e dissipada por custos de transação.\n"
     "\nStatus: amplamente suportada para horizontes de médio/longo prazo.",
     "Fama (1965, 1988)"),
    ("Forma Semi-forte", NAVY,
     "Preços refletem toda informação publicamente disponível: "
     "balanços, releases de resultados, dados macro, notícias e análises.\n"
     "\nImplicação: análise fundamentalista não supera o mercado de forma consistente "
     "após ajuste por risco.\n"
     "\nEvidência: event studies mostram ajuste rápido, mas anomalias como "
     "value e momentum contradizem a forma semi-forte.\n"
     "\nStatus: parcialmente suportada — anomalias são documentadas.",
     "Fama & French (1993); Ball & Brown (1968)"),
    ("Forma Forte", RGBColor(0xB0, 0x33, 0x2A),
     "Preços refletem até informação privada (insider information) "
     "— gestores com acesso privilegiado também não gerariam alfa.\n"
     "\nImplicação: nenhum agente, em nenhuma circunstância, "
     "bate o mercado de forma consistente.\n"
     "\nEvidência: REJEITADA empiricamente. Insider trading gera "
     "retorno anormal antes de anúncios relevantes.\n"
     "\nStatus: não suportada — reguladores proíbem insider trading "
     "exatamente porque ele funciona.",
     "Seyhun (1986); Jeng, Metrick & Zeckhauser (2003)"),
]
for i, (title, color, body, cit) in enumerate(forms):
    rect(s, 0.4 + i * 4.18, 2.08, 3.98, 4.72, fill=LIGHT)
    rect(s, 0.4 + i * 4.18, 2.08, 3.98, 0.42, fill=color)
    txt(s, 0.6 + i * 4.18, 2.13, 3.78, 0.35, title, size=13, bold=True, color=WHITE)
    txt(s, 0.6 + i * 4.18, 2.58, 3.78, 3.5, body, size=10.5, color=NAVY)
    txt(s, 0.6 + i * 4.18, 6.08, 3.78, 0.52, cit, size=9, italic=True, color=GRAY)

ref(s, "Fama, E.F. (1970). Efficient Capital Markets: A Review of Theory and Empirical Work. "
       "J. of Finance, 25(2), 383–417")


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — ANOMALIAS DE MERCADO
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Anomalias de Mercado — O que Quebra a HME",
       "Padrões sistemáticos de retorno que persistem após controle de risco e custos")

txt(s, 0.5, 1.3, 12.3, 0.48,
    "Se os mercados fossem perfeitamente eficientes, nenhum padrão sistemático deveria sobreviver "
    "após ser descoberto (a descoberta eliminaria a anomalia). Décadas de pesquisa documentaram "
    "anomalias robustas — fatores que geram prêmio de retorno mesmo em mercados desenvolvidos.",
    size=13, color=NAVY)

rect(s, 0.5, 1.93, 12.3, 0.03, fill=GOLD)

anomalies = [
    ("Size\n(Tamanho)", BLUE,
     "Small caps > large caps.\nBanz (1981) first to document.\nRisco de liquidez e informação assimétrica.\nRetorno extra ~3% a.a. histórico (EUA).",
     "Banz (1981)\nFama & French (1993)"),
    ("Value\n(Valor)", GREEN,
     "Ações baratas (P/B, P/L baixo) > caras.\nGraham & Dodd (1934) prática.\nFama & French (1993) formalização.\nRisco de distress financeiro.",
     "Fama & French (1993)\nLakonishok et al. (1994)"),
    ("Momentum\n★ Nossa estratégia", GOLD,
     "Vencedores recentes continuam vencendo.\nSinal 12-1 meses é o canônico.\nPersiste em 40+ países e asset classes.\nUnderreaction comportamental.",
     "Jegadeesh & Titman (1993)\nCarhart (1997)"),
    ("Quality\n(Qualidade)", RGBColor(0x6C, 0x35, 0x83),
     "Empresas lucrativas e estáveis geram mais.\nMercado subprecifica persistência de lucros.\nFator QMJ: Quality Minus Junk.\nNovy-Marx (2013) formalização.",
     "Novy-Marx (2013)\nFama & French (2015)"),
    ("Low Vol\n(Paradoxo)", RGBColor(0xBF, 0x6A, 0x02),
     "Ações menos voláteis > mais voláteis.\nVIOLA a relação risco-retorno do CAPM.\nExplicação: restrições de alavancagem.\nFrazzini & Pedersen (2014).",
     "Black (1972)\nFrazzini & Pedersen (2014)"),
]

for i, (name, color, desc, cit) in enumerate(anomalies):
    x = 0.32 + i * 2.58
    rect(s, x, 2.05, 2.48, 4.72, fill=LIGHT)
    rect(s, x, 2.05, 2.48, 0.55, fill=color)
    txt(s, x + 0.1, 2.08, 2.28, 0.5, name,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, 2.67, 2.28, 2.8, desc, size=11, color=NAVY)
    txt(s, x + 0.1, 5.6, 2.28, 0.9, cit, size=9, italic=True, color=GRAY)

ref(s, "Fama & French (1993). Common Risk Factors in Returns on Stocks and Bonds. "
       "J. of Financial Economics, 33(1), 3–56")


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — FACTOR INVESTING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Factor Investing — Do Stock-Picking a Fatores Sistemáticos",
       "A evolução da gestão ativa: substituir julgamento subjetivo por regras replicáveis e auditáveis")

txt(s, 0.5, 1.3, 12.3, 0.52,
    "Factor investing (smart beta) constrói portfólios expostos a fatores de risco documentados de forma "
    "sistemática e transparente — capturando prêmios de retorno sem depender de previsão de mercado "
    "ou habilidades individuais de stock-picking. É a base da gestão quantitativa moderna.",
    size=13, color=NAVY)

rect(s, 0.5, 1.97, 12.3, 0.03, fill=GOLD)

# Linha do tempo de modelos
txt(s, 0.5, 2.05, 12.3, 0.35,
    "Evolução dos modelos de fatores:", size=13, bold=True, color=NAVY)
eras = [
    ("1960s–1970s", "CAPM\n1 fator: mercado", BLUE),
    ("1990s",       "Fama-French\n3 fatores: mercado + size + value", GREEN),
    ("1997",        "Carhart\n4 fatores: + momentum", RGBColor(0x6C, 0x35, 0x83)),
    ("2015",        "Fama-French 5\n5 fatores: + profitability + investment", GOLD),
    ("2010s–hoje",  "Multi-fator + ML\nAlternative data, nowcasting", RGBColor(0xBF, 0x6A, 0x02)),
]
for i, (era, desc, color) in enumerate(eras):
    rect(s, 0.4 + i * 2.55, 2.45, 2.45, 1.85, fill=color)
    txt(s, 0.55 + i * 2.55, 2.52, 2.2, 0.38,
        era, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 0.55 + i * 2.55, 2.92, 2.2, 1.3,
        desc, size=11, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, 0.5, 4.42, 12.3, 0.35,
    "Por que fatores funcionam? Três hipóteses complementares:", size=13, bold=True, color=NAVY)
bullets(s, 0.5, 4.82, 6.2, 2.0,
        ["Risco compensado: o fator representa risco econômico real (ex.: distress, liquidez, duration)",
         "Comportamento: investidores sistematicamente erram por vieses cognitivos (overconfidence, ancoragem)",
         "Fricções estruturais: restrições institucionais (alavancagem, benchmark) criam oportunidades persistentes"],
        size=12, color=NAVY, spacing=7)

rect(s, 6.9, 4.42, 6.0, 2.4, fill=LIGHT)
txt(s, 7.1, 4.52, 5.7, 0.35, "Nosso Posicionamento:", size=13, bold=True, color=NAVY)
txt(s, 7.1, 4.92, 5.6, 1.75,
    "Exploramos o fator momentum no IBOVESPA. A estratégia é sistemática, "
    "baseada em regra clara (sinal 12-1), backtestável com dados públicos "
    "e defensável com referências acadêmicas primárias — exatamente o que "
    "o júri do Desafio Itaú avalia.",
    size=12, color=NAVY)

ref(s, "Ang, A. (2014). Asset Management: A Systematic Approach to Factor Investing. Oxford UP; "
       "Carhart, M. (1997). On Persistence in Mutual Fund Performance. J. of Finance, 52(1), 57–82")


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — O EFEITO MOMENTUM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "O Efeito Momentum",
       "O prêmio de risco mais robusto e mais controverso das finanças quantitativas")

txt(s, 0.5, 1.3, 12.3, 0.52,
    "Jegadeesh & Titman (1993) demonstraram que comprar as ações vencedoras dos últimos 12 meses "
    "e vender as perdedoras gera retorno anormal significativo nos próximos 3–12 meses. "
    "O resultado foi replicado em mais de 40 países e múltiplas classes de ativos.",
    size=13, color=NAVY)

rect(s, 0.5, 1.97, 12.3, 0.03, fill=GOLD)

quadrants = [
    ("Definição Formal", BLUE,
     "Cross-sectional momentum: em cada data de rebalanceamento, ranquear todos os ativos "
     "do universo por seu retorno acumulado nos últimos J meses (excluindo o mais recente — 'skip-1 month'). "
     "Comprar os N melhores (winners) e vender/evitar os N piores (losers). "
     "Nossa implementação: J=11 meses de sinal + 1 mês de skip = sinal '12-1'.",
     "Jegadeesh & Titman (1993)"),
    ("Parâmetros Canônicos", GREEN,
     "J = 3 a 12 meses de formação do sinal (mais comum: 12-1)\n"
     "K = 1 a 6 meses de holding period (rebalanceamento)\n"
     "Skip: 1 mês (evita efeito de reversão de curto prazo / microestrutura)\n"
     "Retorno bruto histórico nos EUA: ~1,0% a.m. (12% a.a.) — antes dos custos",
     "Jegadeesh & Titman (1993, 2001)"),
    ("Robustez Global", RGBColor(0x6C, 0x35, 0x83),
     "Rouwenhorst (1998): confirmado em 12 países europeus, mesmo padrão\n"
     "Asness, Moskowitz & Pedersen (2013): funciona em ações, bonds, câmbio, commodities\n"
     "Griffin et al. (2003): presente em todos os continentes incluindo mercados emergentes\n"
     "Hou et al. (2023): sobrevive à replication crisis — p-value < 0.001",
     "Asness, Moskowitz & Pedersen (2013)"),
    ("Time-series vs Cross-sectional", RGBColor(0xBF, 0x6A, 0x02),
     "Time-series (absoluto): comprar ativo que gerou retorno positivo no período, "
     "vender/sair do ativo que gerou retorno negativo — cada ativo vs. si mesmo.\n"
     "Cross-sectional (relativo): ranquear ativos entre si — comprar melhores, "
     "vender piores do universo. Menos sensível ao nível geral do mercado.\n"
     "Nossa escolha: cross-sectional (mais robusto e mais citado).",
     "Moskowitz, Ooi & Pedersen (2012)"),
]

for i, (title, color, body, cit) in enumerate(quadrants):
    row, col = i // 2, i % 2
    x = 0.4 + col * 6.45
    y = 2.1 + row * 2.65
    rect(s, x, y, 6.2, 2.5, fill=LIGHT)
    rect(s, x, y, 6.2, 0.4, fill=color)
    txt(s, x + 0.15, y + 0.06, 5.9, 0.32, title, size=13, bold=True, color=WHITE)
    txt(s, x + 0.15, y + 0.5, 5.9, 1.55, body, size=11, color=NAVY)
    txt(s, x + 0.15, y + 2.12, 5.9, 0.3, cit, size=9, italic=True, color=GRAY)

ref(s, "Jegadeesh, N. & Titman, S. (1993). Returns to Buying Winners and Selling Losers. "
       "J. of Finance, 48(1), 65–91. | Asness, Moskowitz & Pedersen (2013). J. of Finance, 68(3), 929–985")


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — BEHAVIORAL FINANCE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Por que o Momentum Persiste? — Behavioral Finance",
       "Três mecanismos comportamentais que impedem a arbitragem imediata do prêmio")

txt(s, 0.5, 1.3, 12.3, 0.48,
    "Se os mercados fossem perfeitamente racionais, o momentum desapareceria assim que fosse descoberto "
    "(arbitragem eliminaria o prêmio). A literatura de finanças comportamentais explica sua persistência "
    "por vieses cognitivos sistemáticos que não são arbitrável sem custo e risco.",
    size=13, color=NAVY)

rect(s, 0.5, 1.93, 12.3, 0.03, fill=GOLD)

mechs = [
    ("Underreaction a Notícias", RGBColor(0xC0, 0x39, 0x2B),
     "Mecanismo: investidores atualizam suas crenças muito lentamente "
     "ao receber novas informações relevantes (ancoragem + conservadorismo cognitivo).\n\n"
     "Resultado: boas notícias são gradualmente incorporadas ao preço ao longo de meses, "
     "gerando tendência de preço persistente mesmo após o anúncio.\n\n"
     "Modelo formal: Hong & Stein (1999) propõem dois grupos — newswatchers "
     "(processam fundamental) e momentum traders (seguem preço). Nenhum processa "
     "toda a informação → momentum emerge no equilíbrio.\n\n"
     "Evidência: post-earnings drift — preços continuam a subir por meses após "
     "surpresa positiva nos resultados (Ball & Brown, 1968).",
     "Hong & Stein (1999); Barberis, Shleifer & Vishny (1998)"),
    ("Efeito Disposição", RGBColor(0xE6, 0x8A, 0x00),
     "Mecanismo: investidores têm aversão assimétrica a perdas (prospect theory "
     "de Kahneman & Tversky, 1979). Tendem a vender ganhadores cedo demais "
     "('travar lucro') e segurar perdedores por tempo demais ('esperando recuperar').\n\n"
     "Resultado: vencedores ficam subprecificados por mais tempo do que deveriam "
     "(oferta precoce os deprime), perdedores ficam sobreprecificados "
     "(demanda artificial os sustenta) → momentum no cross-section.\n\n"
     "Modelo formal: Grinblatt & Han (2005) mostram que disposition effect "
     "gera autocorrelação de retornos em equilíbrio.\n\n"
     "Evidência: ações com maior ganho de capital latente têm momentum mais forte.",
     "Shefrin & Statman (1985); Grinblatt & Han (2005)"),
    ("Herding e Momentum Cascade", RGBColor(0x16, 0x80, 0x6A),
     "Mecanismo: gestores institucionais imitam uns aos outros por risco reputacional "
     "— desviar do consenso é arriscado para a carreira mesmo que seja a decisão correta. "
     "Resultado: capital flui para vencedores recentes de forma auto-reforçada.\n\n"
     "Resultado: tendências de preço se amplificam além do fundamento econômico, "
     "criando overshooting que eventualmente reverte — explicando a reversão de longo prazo "
     "(De Bondt & Thaler, 1985) que coexiste com momentum de médio prazo.\n\n"
     "Modelo formal: Scharfstein & Stein (1990) — herding pode ser racional para "
     "o gestor individualmente mesmo sendo irracional para o mercado agregado.\n\n"
     "Evidência: gestores de fundos compram ações que já subiram (Grinblatt et al., 1995).",
     "Scharfstein & Stein (1990); Daniel, Hirshleifer & Subrahmanyam (1998)"),
]

for i, (title, color, body, cit) in enumerate(mechs):
    rect(s, 0.35 + i * 4.32, 2.05, 4.15, 4.72, fill=LIGHT)
    rect(s, 0.35 + i * 4.32, 2.05, 4.15, 0.42, fill=color)
    txt(s, 0.55 + i * 4.32, 2.1, 3.95, 0.35, title, size=12, bold=True, color=WHITE)
    txt(s, 0.55 + i * 4.32, 2.55, 3.95, 3.2, body, size=10, color=NAVY)
    txt(s, 0.55 + i * 4.32, 5.83, 3.95, 0.72, cit, size=9, italic=True, color=GRAY)

ref(s, "Hong, H. & Stein, J. (1999). J. of Finance, 54(6), 2143–2184; "
       "Kahneman, D. & Tversky, A. (1979). Prospect Theory. Econometrica, 47(2), 263–291")


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — MOMENTUM NO BRASIL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Momentum no Brasil — Evidência Empírica",
       "O efeito existe no IBOVESPA? Particularidades do mercado brasileiro")

txt(s, 0.5, 1.3, 12.3, 0.48,
    "A literatura internacional é robusta, mas particularidades do mercado brasileiro — menor número "
    "de ações, alta volatilidade macroeconômica, turnover do índice, e custo de transação mais alto — "
    "levantam questões sobre a transposição direta dos resultados dos EUA.",
    size=13, color=NAVY)

rect(s, 0.5, 1.93, 12.3, 0.03, fill=GOLD)

rect(s, 0.4, 2.05, 7.9, 4.72, fill=LIGHT)
txt(s, 0.6, 2.12, 7.6, 0.38, "O que a pesquisa acadêmica brasileira encontrou:",
    size=14, bold=True, color=NAVY)

br_findings = [
    "Mussa et al. (2012): momentum de 6 e 12 meses é estatisticamente significativo no "
    "mercado acionário brasileiro para o período 1995–2008 (BOVESPA)",
    "Moreira & Mussa (2010): estratégia 12-1 gera ~1,5% a.m. bruto de custos, com "
    "concentração de performance em períodos de baixa volatilidade",
    "Croteu et al. (2019): momentum mais fraco (às vezes reverso) em períodos de "
    "crise macroeconômica (2002, 2008, 2014–2016) — risco de crash é real no Brasil",
    "Conclusão consensual: o efeito é real mas menos estável que nos EUA — "
    "gestão de risco rigorosa e walk-forward são obrigatórios",
    "Amostra pequena (~70–80 ações) exige cautela estatística: p-values são inflados "
    "por data snooping (múltiplos testes — nosso tema da Aula 8)",
]
bullets(s, 0.6, 2.58, 7.6, 3.9, br_findings, size=12, color=NAVY, spacing=7)

rect(s, 8.6, 2.05, 4.4, 4.72, fill=NAVY)
txt(s, 8.8, 2.12, 4.1, 0.38, "Particularidades do Brasil:",
    size=13, bold=True, color=GOLD)
br_parts = [
    "~80 ações (vs. ~3.000 nos EUA)",
    "IBOVESPA reconstitui a cada 4 meses",
    "Muitas ações com spread bid-ask alto",
    "Selic e câmbio afetam mais que nos EUA",
    "Custo: corretagem + B3 fee + IR (15%)",
    "Custo de aluguel para posições vendidas",
    "Implicação: backtest com custos reais é obrigatório — não opcional",
]
bullets(s, 8.8, 2.6, 4.0, 3.9, br_parts, size=12, color=LIGHT, spacing=7)

ref(s, "Mussa, A. et al. (2012). A Estratégia de Momento no Mercado Acionário Brasileiro. "
       "Rev. de Contabilidade & Finanças, 23(58), 42–53; Rouwenhorst (1998). SSRN Working Paper")


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — O DESAFIO ITAÚ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "O Desafio Quant AI — Itaú Asset Management",
       "Contexto da competição, critérios de avaliação e o que diferencia equipes ganhadoras")

rect(s, 0.4, 1.3, 6.0, 5.97, fill=LIGHT)
txt(s, 0.6, 1.38, 5.7, 0.38, "O Desafio:", size=14, bold=True, color=NAVY)
left_items = [
    "Competição universitária patrocinada pela Itaú Asset Management, "
    "uma das maiores gestoras do Brasil (>R$ 700 bi sob gestão)",
    "Participantes: equipes de graduação e pós-graduação brasileiras, "
    "avaliadas por profissionais sênior da indústria",
    "Objetivo: desenvolver uma estratégia quantitativa completa — "
    "da hipótese econômica ao relatório técnico final",
    "Entregável: relatório técnico escrito + apresentação oral de 20 min + Q&A com o júri",
    "O que diferencia: rigor metodológico e capacidade de defender "
    "cada escolha com fundamento econômico e evidência estatística",
]
bullets(s, 0.6, 1.85, 5.7, 4.7, left_items, size=12, color=NAVY, spacing=9)

rect(s, 6.85, 1.3, 6.1, 5.97, fill=NAVY)
txt(s, 7.05, 1.38, 5.8, 0.38, "Critérios de Avaliação:", size=14, bold=True, color=GOLD)
criteria = [
    ("Tese de Investimento", "20%",
     "Clareza da hipótese e fundamento econômico"),
    ("Dados & Metodologia", "20%",
     "Qualidade, limpeza, ausência de look-ahead bias"),
    ("Rigor do Backtest",   "20%",
     "Walk-forward, custos reais, múltiplos cenários"),
    ("Resultados",          "15%",
     "Sharpe, Sortino, drawdown, Calmar, consistência"),
    ("Análise de Risco",    "15%",
     "Sensibilidade a parâmetros, cenários de stress"),
    ("Comunicação",         "10%",
     "Relatório claro, defesa oral objetiva e convincente"),
]
for i, (crit, pct, desc) in enumerate(criteria):
    y = 1.88 + i * 0.84
    rect(s, 6.95, y, 5.9, 0.76, fill=RGBColor(0x1A, 0x2E, 0x5C))
    txt(s, 7.1, y + 0.07, 4.2, 0.3, crit, size=12, bold=True, color=WHITE)
    txt(s, 7.1, y + 0.4, 4.2, 0.3, desc, size=10, italic=True, color=LIGHT)
    txt(s, 11.5, y + 0.15, 1.2, 0.42, pct, size=20, bold=True,
        color=GOLD, align=PP_ALIGN.RIGHT)

ref(s, "Itaú Asset Management. Regulamento do Desafio Quant AI (2025). itauassetmanagement.com.br")


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — NOSSA ESTRATÉGIA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Nossa Estratégia — Momentum Cross-Sectional no IBOVESPA",
       "Definição completa: universo, sinal, portfólio, rebalanceamento e validação")

specs = [
    ("Universo de Ativos", BLUE,
     "~77 ações constituintes do IBOVESPA\n"
     "Fonte: Yahoo Finance via yfinance (.SA)\n"
     "Período: 2014–2024 (10 anos de backtest)\n"
     "Filtro: ações com dados suficientes em cada data"),
    ("Sinal (Fator)", GOLD,
     "Retorno acumulado 11 meses excluindo o mais recente\n"
     "Python: ret_mensais.shift(2).rolling(11).sum()\n"
     "Skip-1: evita reversão de microestrutura de curto prazo\n"
     "Versão v2: vol-ajustado (÷ vol_rolling_63d)"),
    ("Portfólio", GREEN,
     "Top 10 ações por sinal → long-only\n"
     "v1: equal-weight (1/N) — simples e robusto\n"
     "v2: pesos proporcionais ao sinal / vol\n"
     "Sem posições vendidas (long-only = operacional)"),
    ("Rebalanceamento & Custos", RGBColor(0x6C, 0x35, 0x83),
     "Frequência: mensal (último dia útil)\n"
     "Hipótese: execução no fechamento seguinte\n"
     "Custos modelados: 0%, 0,3%, 0,5%, 1,0%/turno\n"
     "Break-even análisis incluso no backtest"),
    ("Gestão de Risco & Validação", RGBColor(0xBF, 0x6A, 0x02),
     "Walk-forward: 48m treino / 12m teste (rolling)\n"
     "Deflated Sharpe Ratio (Bailey & LdP, 2014)\n"
     "Análise de sensibilidade: J = 3 a 36 meses\n"
     "Múltiplos cenários de custo e N de ações"),
]

for i, (title, color, body) in enumerate(specs):
    row = i // 3
    col = i % 3
    if row == 0:
        x, y, w2, h2 = 0.35 + col * 4.25, 1.32, 4.05, 2.85
    else:
        x, y, w2, h2 = 0.35 + (i - 3) * 6.38, 4.3, 6.12, 2.95
    rect(s, x, y, w2, h2, fill=LIGHT)
    rect(s, x, y, w2, 0.42, fill=color)
    txt(s, x + 0.15, y + 0.07, w2 - 0.2, 0.34, title,
        size=12, bold=True, color=WHITE)
    txt(s, x + 0.15, y + 0.52, w2 - 0.2, h2 - 0.65, body, size=12, color=NAVY)

ref(s, "Jegadeesh & Titman (1993); Bailey & López de Prado (2014). "
       "The Deflated Sharpe Ratio. Financial Analysts Journal, 70(5), 94–107")


# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — ROADMAP DO INTENSIVÃO
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Roadmap do Intensivão Quant AI — 10 Aulas",
       "Da coleta de dados ao relatório final: uma estratégia quant completa em 5 semanas")

aulas = [
    ("01", "Kickoff",           "Fundamentos de mercado, momentum, estrutura do projeto"),
    ("02", "Dados",             "yfinance, IBOVESPA, limpeza, parquet, pipeline de dados"),
    ("03", "EDA",               "Distribuição de retornos, correlações, estacionariedade"),
    ("04", "Sinal v1",          "Momentum 12-1, implementação, diagnóstico do sinal"),
    ("05", "Backtest v1",       "Portfólio long-only, métricas (Sharpe, MDD, Calmar)"),
    ("06", "Alocação",          "Equal-weight vs vol-weight vs otimização (MV)"),
    ("07", "Sinal v2",          "Momentum vol-ajustado, comparação de versões"),
    ("08", "Backtest Rigoroso", "Look-ahead bias, overfitting, DSR, custos, walk-forward"),
    ("09", "GenAI + Análise",   "API Anthropic, prompt engineering, rascunho do relatório"),
    ("10", "Relatório+Defesa",  "Tear sheet, rubrica de banca, simulação de defesa"),
]

colors_row = [BLUE, RGBColor(0x6C, 0x35, 0x83)]
for i, (num, titulo, desc) in enumerate(aulas):
    row, col = i // 5, i % 5
    x = 0.32 + col * 2.56
    y = 1.35 + row * 2.93
    is_today = (i == 0)
    bg_card = RGBColor(0xFF, 0xF0, 0xC8) if is_today else LIGHT
    rect(s, x, y, 2.42, 2.72, fill=bg_card)
    num_color = GOLD if is_today else colors_row[row]
    rect(s, x, y, 2.42, 0.5, fill=num_color if not is_today else NAVY)
    txt(s, x + 0.1, y + 0.08, 1.5, 0.38,
        f"Aula {num}", size=13, bold=True,
        color=GOLD if not is_today else WHITE)
    if is_today:
        txt(s, x + 1.6, y + 0.1, 0.75, 0.32,
            "HOJE", size=9, bold=True,
            color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, y + 0.6, 2.22, 0.42,
        titulo, size=12, bold=True, color=NAVY)
    txt(s, x + 0.1, y + 1.05, 2.22, 1.55,
        desc, size=10, color=NAVY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — REFERÊNCIAS BIBLIOGRÁFICAS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=WHITE)
header(s, "Referências Bibliográficas",
       "Leitura recomendada — todos disponíveis no Google Scholar ou SSRN (acesso aberto)")

refs_list = [
    "Markowitz, H. (1952). Portfolio Selection. The Journal of Finance, 7(1), 77–91.",
    "Sharpe, W.F. (1964). Capital Asset Prices: A Theory of Market Equilibrium Under Conditions of Risk. J. of Finance, 19(3), 425–442.",
    "Black, F. & Scholes, M. (1973). The Pricing of Options and Corporate Liabilities. Journal of Political Economy, 81(3), 637–654.",
    "Fama, E.F. (1970). Efficient Capital Markets: A Review of Theory and Empirical Work. Journal of Finance, 25(2), 383–417.",
    "Banz, R.W. (1981). The Relationship Between Return and Market Value of Common Stocks. Journal of Financial Economics, 9(1), 3–18.",
    "Mehra, R. & Prescott, E.C. (1985). The Equity Premium: A Puzzle. Journal of Monetary Economics, 15(2), 145–161.",
    "Kahneman, D. & Tversky, A. (1979). Prospect Theory: An Analysis of Decision Under Risk. Econometrica, 47(2), 263–291.",
    "Fama, E.F. & French, K.R. (1993). Common Risk Factors in the Returns on Stocks and Bonds. Journal of Financial Economics, 33(1), 3–56.",
    "Jegadeesh, N. & Titman, S. (1993). Returns to Buying Winners and Selling Losers: Implications for Stock Market Efficiency. Journal of Finance, 48(1), 65–91.",
    "Carhart, M.M. (1997). On Persistence in Mutual Fund Performance. Journal of Finance, 52(1), 57–82.",
    "Barberis, N., Shleifer, A. & Vishny, R. (1998). A Model of Investor Sentiment. Journal of Financial Economics, 49(3), 307–343.",
    "Hong, H. & Stein, J.C. (1999). A Unified Theory of Underreaction, Momentum Trading, and Overreaction. Journal of Finance, 54(6), 2143–2184.",
    "Rouwenhorst, K.G. (1998). International Momentum Strategies. Journal of Finance, 53(1), 267–284.",
    "DeMiguel, V., Garlappi, L. & Uppal, R. (2009). Optimal Versus Naive Diversification. Review of Financial Studies, 22(5), 1915–1953.",
    "Asness, C., Moskowitz, T. & Pedersen, L.H. (2013). Value and Momentum Everywhere. Journal of Finance, 68(3), 929–985.",
    "Ang, A. (2014). Asset Management: A Systematic Approach to Factor Investing. Oxford University Press.",
    "Bailey, D.H. & López de Prado, M. (2014). The Deflated Sharpe Ratio: Correcting for Selection Bias, Backtest Overfitting and Non-Normality. Financial Analysts Journal, 70(5), 94–107.",
    "Mussa, A. et al. (2012). A Estratégia de Momento no Mercado Acionário Brasileiro. Revista de Contabilidade & Finanças, 23(58), 42–53.",
    "Dimson, E., Marsh, P. & Staunton, M. (2002). Triumph of the Optimists: 101 Years of Global Investment Returns. Princeton University Press.",
    "Fabozzi, F.J. (2012). The Handbook of Fixed Income Securities, 8ª edição. McGraw-Hill.",
]

mid = len(refs_list) // 2
tb_left = s.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.9))
tf_l = tb_left.text_frame
tf_l.word_wrap = True
for i, r in enumerate(refs_list[:mid]):
    p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run()
    run.text = f"{i+1}. {r}"
    run.font.size = Pt(9.5)
    run.font.color.rgb = NAVY

tb_right = s.shapes.add_textbox(Inches(6.85), Inches(1.3), Inches(6.2), Inches(5.9))
tf_r = tb_right.text_frame
tf_r.word_wrap = True
for i, r in enumerate(refs_list[mid:]):
    p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run()
    run.text = f"{mid+i+1}. {r}"
    run.font.size = Pt(9.5)
    run.font.color.rgb = NAVY


# ── Salvar ───────────────────────────────────────────────────
output = "slides-aula-01-kickoff.pptx"
prs.save(output)
print(f"OK: {output} salvo com sucesso!")
print(f"  Total de slides: {len(prs.slides)}")
print(f"  Formato: {prs.slide_width.inches:.2f}'' x {prs.slide_height.inches:.2f}'' (16:9)")
